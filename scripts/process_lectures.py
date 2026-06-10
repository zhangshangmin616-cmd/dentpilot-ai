from __future__ import annotations

import json
import re
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

PROJECT_ROOT = Path(__file__).resolve().parents[1]
RAW_ROOT = PROJECT_ROOT / "data" / "raw_lectures"
PROCESSED_TEXT_ROOT = PROJECT_ROOT / "data" / "processed_lectures"
KNOWLEDGE_ROOT = PROJECT_ROOT / "data" / "processed_knowledge_base"
LOG_ROOT = PROJECT_ROOT / "data" / "lecture_processing_logs"
LOG_PATH = LOG_ROOT / "process_log.json"
MAX_MARKDOWN_CHARS = 250_000
SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".txt", ".docx"}
LEGACY_PPT_WARNING = "PPT legacy format detected. Please convert to PPTX or PDF."


def normalize_whitespace(text: str) -> str:
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def safe_name(value: str) -> str:
    cleaned = re.sub(r"[^\w.\-]+", "_", value.strip(), flags=re.UNICODE)
    return cleaned.strip("_") or "untitled"


def extract_pdf(path: Path) -> tuple[str, dict[str, Any]]:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise RuntimeError("Missing dependency pypdf. Run: pip install pypdf") from exc

    reader = PdfReader(str(path))
    page_texts: list[str] = []
    pages_with_text = 0

    for index, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text() or ""
        except Exception as exc:  # Keep other pages if one extraction fails.
            text = f"[Page {index} extraction failed: {exc}]"
        if text.strip():
            pages_with_text += 1
        page_texts.append(f"\n\n## Page {index}\n\n{text.strip()}")

    meta = {
        "type": "pdf",
        "total_pages": len(reader.pages),
        "pages_with_text": pages_with_text,
    }
    return normalize_whitespace("\n".join(page_texts)), meta


def extract_pptx(path: Path) -> tuple[str, dict[str, Any]]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("Missing dependency python-pptx. Run: pip install python-pptx") from exc

    presentation = Presentation(str(path))
    slide_texts: list[str] = []
    slides_with_text = 0

    for index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    parts.append(" | ".join(cell.text for cell in row.cells))
        text = normalize_whitespace("\n".join(parts))
        if text:
            slides_with_text += 1
        slide_texts.append(f"\n\n## Slide {index}\n\n{text}")

    meta = {
        "type": "pptx",
        "total_slides": len(presentation.slides),
        "slides_with_text": slides_with_text,
    }
    return normalize_whitespace("\n".join(slide_texts)), meta


def extract_docx(path: Path) -> tuple[str, dict[str, Any]]:
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError("Missing dependency python-docx. Run: pip install python-docx") from exc

    document = Document(str(path))
    parts: list[str] = []

    for paragraph in document.paragraphs:
        if paragraph.text.strip():
            parts.append(paragraph.text)

    for table in document.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text.strip() for cell in row.cells if cell.text.strip()))

    meta = {
        "type": "docx",
        "paragraphs": len(document.paragraphs),
        "tables": len(document.tables),
    }
    return normalize_whitespace("\n".join(parts)), meta


def extract_txt(path: Path) -> tuple[str, dict[str, Any]]:
    for encoding in ("utf-8", "utf-8-sig", "gb18030", "latin-1"):
        try:
            return normalize_whitespace(path.read_text(encoding=encoding)), {
                "type": "txt",
                "encoding": encoding,
            }
        except UnicodeDecodeError:
            continue
    return normalize_whitespace(path.read_text(errors="ignore")), {
        "type": "txt",
        "encoding": "unknown_fallback",
    }


def extract_text(path: Path) -> tuple[str, dict[str, Any]]:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return extract_pdf(path)
    if suffix == ".pptx":
        return extract_pptx(path)
    if suffix == ".docx":
        return extract_docx(path)
    if suffix == ".txt":
        return extract_txt(path)
    raise ValueError(f"Unsupported file type: {suffix}")


def split_into_sections(text: str) -> list[tuple[str, str]]:
    heading_pattern = re.compile(
        r"(?m)^(?:#{1,4}\s+.+|\d{1,2}[.)、]\s+.+|Slide\s+\d+|Page\s+\d+)$"
    )
    matches = list(heading_pattern.finditer(text))
    if not matches:
        return [("Full lecture text", text)]

    sections: list[tuple[str, str]] = []
    for index, match in enumerate(matches):
        start = match.start()
        end = matches[index + 1].start() if index + 1 < len(matches) else len(text)
        chunk = text[start:end].strip()
        title = match.group(0).strip("# ").strip()
        if chunk:
            sections.append((title, chunk))
    return sections or [("Full lecture text", text)]


def important_terms(text: str, limit: int = 30) -> list[str]:
    terms = re.findall(r"\b[A-Za-z][A-Za-z][A-Za-z\-]{2,}\b", text.lower())
    stop = {
        "this",
        "that",
        "with",
        "from",
        "have",
        "were",
        "been",
        "their",
        "there",
        "which",
        "lecture",
        "slide",
        "page",
        "dental",
    }
    counts: dict[str, int] = {}
    for term in terms:
        if term in stop:
            continue
        counts[term] = counts.get(term, 0) + 1
    return [term for term, _ in sorted(counts.items(), key=lambda item: (-item[1], item[0]))[:limit]]


def compact_section(title: str, text: str) -> str:
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    content_lines = [
        line
        for line in lines
        if len(line) >= 8 and not re.match(r"^(page|slide)\s+\d+$", line, flags=re.I)
    ]
    preview = content_lines[:12]
    terms = important_terms(text, 12)

    return "\n".join(
        [
            f"### {title}",
            "",
            "**High-yield lecture notes**",
            *(f"- {line[:420]}" for line in preview),
            "",
            "**Likely exam focus**",
            "- Definition, classification, mechanism, clinical relevance, diagnosis, treatment principles, or material properties depending on the lecture content.",
            "- Oral exam answers should connect the concept to dental clinical practice where possible.",
            "",
            "**Key terms**",
            *(f"- {term}" for term in terms),
            "",
        ]
    )


def build_raw_markdown(subject: str, extracted_files: list[dict[str, Any]]) -> str:
    parts = [
        f"# Lecture Raw Merge: {subject}",
        "",
        "This file is generated from local lecture files for ElevenLabs Knowledge Base ingestion.",
        "",
    ]
    for item in extracted_files:
        parts.extend(
            [
                f"## Source: {item['filename']}",
                "",
                f"- Characters: {item['characters']}",
                f"- Metadata: `{json.dumps(item['metadata'], ensure_ascii=False)}`",
                "",
                item["text"],
                "",
            ]
        )
    return normalize_whitespace("\n".join(parts)) + "\n"


def build_exam_knowledge_markdown(subject: str, extracted_files: list[dict[str, Any]]) -> str:
    parts = [
        f"# Lecture Exam Knowledge: {subject}",
        "",
        "Use this compact file as a course-aligned knowledge base for oral exam practice.",
        "",
        "## Source Coverage",
        "",
    ]
    for item in extracted_files:
        parts.append(f"- {item['filename']}: {item['characters']} extracted characters")

    parts.extend(["", "## Exam Knowledge Modules", ""])

    for item in extracted_files:
        parts.extend([f"## Source: {item['filename']}", ""])
        for title, section_text in split_into_sections(item["text"]):
            parts.append(compact_section(title, section_text))

    return normalize_whitespace("\n".join(parts)) + "\n"


def write_split_text(base_path: Path, content: str) -> list[str]:
    base_path.parent.mkdir(parents=True, exist_ok=True)
    if len(content) <= MAX_MARKDOWN_CHARS:
        base_path.write_text(content, encoding="utf-8")
        return [str(base_path.relative_to(PROJECT_ROOT))]

    written: list[str] = []
    stem = base_path.stem
    suffix = base_path.suffix
    for index, start in enumerate(range(0, len(content), MAX_MARKDOWN_CHARS), start=1):
        part_path = base_path.with_name(f"{stem}_part{index}{suffix}")
        part_path.write_text(content[start : start + MAX_MARKDOWN_CHARS], encoding="utf-8")
        written.append(str(part_path.relative_to(PROJECT_ROOT)))
    return written


def process_subject(subject_dir: Path) -> dict[str, Any]:
    subject = subject_dir.name
    processed_subject_dir = PROCESSED_TEXT_ROOT / subject
    processed_subject_dir.mkdir(parents=True, exist_ok=True)

    result: dict[str, Any] = {
        "subject": subject,
        "files": [],
        "warnings": [],
        "errors": [],
        "outputs": [],
    }

    extracted_files: list[dict[str, Any]] = []
    for path in sorted(subject_dir.rglob("*")):
        if not path.is_file() or path.name == ".gitkeep":
            continue

        suffix = path.suffix.lower()
        if suffix == ".ppt":
            warning = f"{path.relative_to(PROJECT_ROOT)}: {LEGACY_PPT_WARNING}"
            result["warnings"].append(warning)
            continue

        if suffix not in SUPPORTED_EXTENSIONS:
            result["warnings"].append(
                f"{path.relative_to(PROJECT_ROOT)}: unsupported file type skipped"
            )
            continue

        try:
            text, metadata = extract_text(path)
            output_txt = processed_subject_dir / f"{safe_name(path.stem)}.txt"
            output_txt.write_text(text + "\n", encoding="utf-8")
            record = {
                "filename": path.name,
                "source_path": str(path.relative_to(PROJECT_ROOT)),
                "processed_text_path": str(output_txt.relative_to(PROJECT_ROOT)),
                "characters": len(text),
                "metadata": metadata,
                "text": text,
            }
            extracted_files.append(record)
            result["files"].append({key: value for key, value in record.items() if key != "text"})
        except Exception as exc:
            result["errors"].append(f"{path.relative_to(PROJECT_ROOT)}: {exc}")

    if extracted_files:
        raw_markdown = build_raw_markdown(subject, extracted_files)
        exam_markdown = build_exam_knowledge_markdown(subject, extracted_files)
        result["outputs"].extend(
            write_split_text(
                KNOWLEDGE_ROOT / f"LECTURE__{subject}__raw_merged.md",
                raw_markdown,
            )
        )
        result["outputs"].extend(
            write_split_text(
                KNOWLEDGE_ROOT / f"LECTURE__{subject}__exam_knowledge.md",
                exam_markdown,
            )
        )

    return result


def build_index(results: list[dict[str, Any]]) -> str:
    parts = [
        "# Lecture Knowledge Base Index",
        "",
        f"Generated at: {datetime.now(timezone.utc).isoformat()}",
        "",
        "Upload the Markdown files in this folder to ElevenLabs Knowledge Base.",
        "",
        "## Subjects",
        "",
    ]
    for result in results:
        parts.extend(
            [
                f"### {result['subject']}",
                f"- Extracted files: {len(result['files'])}",
                f"- Warnings: {len(result['warnings'])}",
                f"- Errors: {len(result['errors'])}",
                "- Outputs:",
                *(f"  - {output}" for output in result["outputs"]),
                "",
            ]
        )
    return "\n".join(parts) + "\n"


def main() -> None:
    RAW_ROOT.mkdir(parents=True, exist_ok=True)
    PROCESSED_TEXT_ROOT.mkdir(parents=True, exist_ok=True)
    KNOWLEDGE_ROOT.mkdir(parents=True, exist_ok=True)
    LOG_ROOT.mkdir(parents=True, exist_ok=True)

    subject_dirs = [path for path in sorted(RAW_ROOT.iterdir()) if path.is_dir()]
    results = [process_subject(subject_dir) for subject_dir in subject_dirs]

    index_path = KNOWLEDGE_ROOT / "LECTURE_INDEX.md"
    index_path.write_text(build_index(results), encoding="utf-8")

    log = {
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "raw_root": str(RAW_ROOT.relative_to(PROJECT_ROOT)),
        "processed_text_root": str(PROCESSED_TEXT_ROOT.relative_to(PROJECT_ROOT)),
        "knowledge_root": str(KNOWLEDGE_ROOT.relative_to(PROJECT_ROOT)),
        "max_markdown_chars": MAX_MARKDOWN_CHARS,
        "subjects": results,
    }
    LOG_PATH.write_text(json.dumps(log, ensure_ascii=False, indent=2), encoding="utf-8")

    print(f"Processed {len(subject_dirs)} subject folders.")
    print(f"Knowledge index: {index_path}")
    print(f"Process log: {LOG_PATH}")


if __name__ == "__main__":
    main()

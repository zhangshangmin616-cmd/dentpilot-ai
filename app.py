import csv
import io
import json
import os
import re
import time
from pathlib import Path
from urllib.parse import quote

import pandas as pd
import requests
import streamlit as st
from streamlit_js_eval import streamlit_js_eval
from docx import Document
from dotenv import load_dotenv
from pypdf import PdfReader
try:
    from pptx import Presentation
except Exception:
    Presentation = None
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.cidfonts import UnicodeCIDFont
from reportlab.pdfbase.ttfonts import TTFont

from ai_engine import generate_study_pack
from clinical_case import (
    ClinicalCaseConfigError,
    ClinicalCaseJSONError,
    generate_clinical_case,
    grade_clinical_case,
)
from oral_exam import (
    OralExamConfigError,
    OralExamJSONError,
    generate_oral_question,
    grade_oral_answer,
)
from question_bank_engine import (
    find_relevant_school_questions,
    generate_written_question_from_bank,
    get_question_bank_status,
    infer_question_focus,
    load_school_question_bank,
)
from realtime_oral_exam import render_realtime_oral_exam_page
from ui_i18n import get_ui_text, normalize_lang
from weakness_analysis import (
    WeaknessAnalysisConfigError,
    WeaknessAnalysisJSONError,
    analyze_weaknesses,
)


load_dotenv(Path(__file__).resolve().parent / ".env", override=False)


SUPABASE_AUTH_URL = "https://nakkcdzpxdggirujgmtk.supabase.co/auth/v1"
DEFAULT_SUPABASE_PUBLISHABLE_KEY = "sb_publishable_mBC1RRvQRbZmNfofqDap2w_z0DjtKrE"
LOCAL_STORAGE_AUTH_KEY = "DENTPILOT_AUTH_SESSION"
LOCAL_STORAGE_UI_LANG_KEY = "DENTPILOT_UI_LANG"
LOCAL_STORAGE_EMPTY = "__DENTPILOT_AUTH_EMPTY__"
AUTH_SESSION_KEYS = (
    "dentpilot_user",
    "dentpilot_access_token",
    "dentpilot_refresh_token",
    "dentpilot_expires_at",
    "auth_user",
    "auth_session",
)

SUBJECT_OPTIONS = [
    ("Dentistry", "口腔医学 Dentistry"),
    ("Endodontics", "牙体牙髓 Endodontics"),
    ("Periodontology", "牙周 Periodontology"),
    ("Oral Surgery", "口腔外科 Oral Surgery"),
    ("Oral Pathology", "口腔病理 Oral Pathology"),
    ("Dental Anatomy", "牙体解剖 Dental Anatomy"),
    ("Pharmacology", "药理学 Pharmacology"),
    ("General Pathology", "普通病理 General Pathology"),
    ("Orthodontics", "口腔正畸 Orthodontics"),
    ("Preventive Dentistry", "口腔预防 Preventive Dentistry"),
    ("Microbiology", "微生物学 Microbiology"),
]

SUBJECT_VALUES = [value for value, _label in SUBJECT_OPTIONS]
SUBJECT_LABELS = dict(SUBJECT_OPTIONS)


def format_subject_option(value: str) -> str:
    if get_ui_lang() == "en":
        return value
    return SUBJECT_LABELS.get(value, value)


def get_ui_lang() -> str:
    return normalize_lang(st.session_state.get("ui_lang", "zh"))


def t(key: str) -> str:
    return get_ui_text(get_ui_lang(), key)


def save_ui_lang_to_local_storage(lang: str) -> None:
    try:
        streamlit_js_eval(
            js_expressions=(
                f"localStorage.setItem({json.dumps(LOCAL_STORAGE_UI_LANG_KEY)}, "
                f"{json.dumps(normalize_lang(lang))}); 'saved';"
            ),
            key=f"ui_lang_save_{int(time.time() * 1000)}",
        )
    except Exception:
        pass


def init_ui_language() -> None:
    st.session_state.setdefault("ui_lang", "zh")
    try:
        query_lang = st.query_params.get("ui_lang") or st.query_params.get("lang")
        if isinstance(query_lang, list):
            query_lang = query_lang[0] if query_lang else ""
        query_lang = normalize_lang(str(query_lang)) if query_lang else ""
        if query_lang in {"zh", "en"}:
            st.session_state["ui_lang"] = query_lang
            st.session_state["ui_lang_loaded"] = True
            return
    except Exception:
        pass
    try:
        raw_value = streamlit_js_eval(
            js_expressions=(
                f"localStorage.getItem({json.dumps(LOCAL_STORAGE_UI_LANG_KEY)}) || "
                f"{json.dumps('')};"
            ),
            key="ui_lang_load",
        )
        if raw_value in {"zh", "en"} and "ui_lang_loaded" not in st.session_state:
            st.session_state["ui_lang"] = raw_value
            st.session_state["ui_lang_loaded"] = True
    except Exception:
        pass


def render_ui_language_selector() -> None:
    lang = get_ui_lang()
    options = ["zh", "en"]
    selected = st.selectbox(
        t("ui_language"),
        options,
        index=options.index(lang),
        format_func=lambda value: get_ui_text(lang, value),
        key="ui_lang_selector",
    )
    selected = normalize_lang(selected)
    if selected != st.session_state.get("ui_lang"):
        st.session_state["ui_lang"] = selected
        save_ui_lang_to_local_storage(selected)
        st.rerun()
    st.caption(t("ui_note"))


def read_config_value(name: str, default: str = "") -> str:
    try:
        secret_value = st.secrets.get(name)
        if secret_value:
            return str(secret_value)
    except Exception:
        pass
    return os.getenv(name, default)


def get_supabase_auth_config() -> tuple[str, str]:
    url = read_config_value(
        "NEXT_PUBLIC_SUPABASE_URL",
        read_config_value("SUPABASE_URL", "https://nakkcdzpxdggirujgmtk.supabase.co"),
    ).rstrip("/")
    key = read_config_value(
        "NEXT_PUBLIC_SUPABASE_ANON_KEY",
        read_config_value("SUPABASE_ANON_KEY", DEFAULT_SUPABASE_PUBLISHABLE_KEY),
    )
    return url, key


def supabase_auth_request(endpoint: str, payload: dict) -> dict:
    supabase_url, publishable_key = get_supabase_auth_config()
    if not supabase_url or not publishable_key:
        raise RuntimeError("Supabase \u767b\u5f55\u5c1a\u672a\u914d\u7f6e\uff0c\u8bf7\u5728 Streamlit Secrets \u4e2d\u6dfb\u52a0 Publishable key\u3002")

    response = requests.post(
        f"{supabase_url}/auth/v1/{endpoint}",
        headers={
            "apikey": publishable_key,
            "Authorization": f"Bearer {publishable_key}",
            "Content-Type": "application/json",
        },
        json=payload,
        timeout=20,
    )
    data = response.json() if response.content else {}
    if response.status_code >= 400:
        message = data.get("msg") or data.get("message") or data.get("error_description") or response.text
        raise RuntimeError(message)
    return data


def sign_in_with_email(email: str, password: str) -> dict:
    return supabase_auth_request(
        "token?grant_type=password",
        {"email": email, "password": password},
    )


def sign_up_with_email(email: str, password: str) -> dict:
    return supabase_auth_request(
        "signup",
        {"email": email, "password": password},
    )


def auth_debug_enabled() -> bool:
    return read_config_value("DENTPILOT_AUTH_DEBUG", "").lower() in {"1", "true", "yes", "on"}


def demo_login_enabled() -> bool:
    return read_config_value("DENTPILOT_DEMO_LOGIN", "").lower() in {"1", "true", "yes", "on"}


def get_demo_user() -> dict:
    return {
        "id": "local-demo-user",
        "email": "demo@dentpilot.local",
    }


def record_auth_debug(**values) -> None:
    if not auth_debug_enabled():
        return
    debug = st.session_state.setdefault("auth_debug", {})
    debug.update(values)


def render_auth_debug() -> None:
    if not auth_debug_enabled():
        return
    debug = st.session_state.get("auth_debug", {})
    st.markdown("#### \u767b\u5f55\u72b6\u6001\u8c03\u8bd5")
    st.caption(
        " | ".join(
            [
                f"session user: {'yes' if debug.get('session_user') else 'no'}",
                f"localStorage auth: {'yes' if debug.get('local_storage_exists') else 'no'}",
                f"restore attempted: {'yes' if debug.get('restore_attempted') else 'no'}",
                f"refresh attempted: {'yes' if debug.get('refresh_attempted') else 'no'}",
                f"current user: {debug.get('current_user_email') or 'none'}",
            ]
        )
    )


def get_auth_headers(access_token: str | None = None) -> dict:
    _, publishable_key = get_supabase_auth_config()
    return {
        "apikey": publishable_key,
        "Authorization": f"Bearer {access_token or publishable_key}",
        "Content-Type": "application/json",
    }


def normalize_supabase_user(user: dict | None) -> dict:
    if not isinstance(user, dict):
        return {}
    metadata = user.get("user_metadata") or {}
    return {
        "id": user.get("id"),
        "email": user.get("email") or metadata.get("email"),
    }


def build_auth_payload(data: dict) -> dict:
    expires_at = data.get("expires_at")
    if not expires_at and data.get("expires_in"):
        expires_at = int(time.time()) + int(data.get("expires_in") or 0)
    return {
        "access_token": data.get("access_token"),
        "refresh_token": data.get("refresh_token"),
        "expires_at": expires_at,
        "user": normalize_supabase_user(data.get("user")),
    }


def safe_int(value, default: int = 0) -> int:
    try:
        return int(float(value))
    except Exception:
        return default


def safe_float(value, default=None):
    try:
        if value is None or value == "":
            return default
        return float(value)
    except Exception:
        return default


def apply_auth_payload(payload: dict, restored: bool = False) -> dict | None:
    if not payload.get("access_token") or not payload.get("refresh_token"):
        return None
    user = normalize_supabase_user(payload.get("user"))
    st.session_state["auth_session"] = payload
    st.session_state["auth_user"] = user
    st.session_state["dentpilot_access_token"] = payload.get("access_token")
    st.session_state["dentpilot_refresh_token"] = payload.get("refresh_token")
    st.session_state["dentpilot_expires_at"] = payload.get("expires_at")
    st.session_state["dentpilot_user"] = user
    record_auth_debug(
        session_user=True,
        current_user_email=user.get("email"),
        restore_attempted=restored,
    )
    if restored and user.get("email"):
        st.session_state["dentpilot_auth_restored_message"] = f"\u5df2\u81ea\u52a8\u767b\u5f55\uff1a{user['email']}"
    return user


def make_local_storage_set_js(payload: dict) -> str:
    storage_key = json.dumps(LOCAL_STORAGE_AUTH_KEY)
    storage_value = json.dumps(json.dumps(payload, ensure_ascii=False, separators=(",", ":")))
    return f"localStorage.setItem({storage_key}, {storage_value}); 'saved';"


def save_auth_session_to_local_storage(session: dict) -> None:
    streamlit_js_eval(
        js_expressions=make_local_storage_set_js(session),
        key=f"auth_local_storage_save_{int(time.time() * 1000)}",
    )
    st.session_state["auth_local_storage_recently_saved"] = True
    record_auth_debug(local_storage_exists=True)


def save_auth_session(data: dict) -> dict | None:
    payload = build_auth_payload(data)
    user = apply_auth_payload(payload)
    if not user:
        return None
    save_auth_session_to_local_storage(payload)
    return user


def load_auth_session_from_local_storage() -> dict | None:
    storage_key = json.dumps(LOCAL_STORAGE_AUTH_KEY)
    empty_value = json.dumps(LOCAL_STORAGE_EMPTY)
    raw_value = streamlit_js_eval(
        js_expressions=f"localStorage.getItem({storage_key}) || {empty_value};",
        key="auth_local_storage_load",
    )

    if raw_value is None:
        attempts = int(st.session_state.get("auth_local_storage_attempts", 0)) + 1
        st.session_state["auth_local_storage_attempts"] = attempts
        st.session_state["auth_local_storage_pending"] = True
        record_auth_debug(local_storage_exists=False, restore_attempted=False)
        return None

    st.session_state["auth_local_storage_pending"] = False
    st.session_state["auth_local_storage_attempts"] = 0
    if raw_value == LOCAL_STORAGE_EMPTY:
        record_auth_debug(local_storage_exists=False, restore_attempted=True)
        return None

    try:
        payload = json.loads(raw_value)
    except Exception:
        clear_auth_session_from_local_storage()
        record_auth_debug(local_storage_exists=True, restore_attempted=True)
        return None

    record_auth_debug(local_storage_exists=True, restore_attempted=True)
    return payload


def clear_auth_session_from_local_storage() -> None:
    storage_key = json.dumps(LOCAL_STORAGE_AUTH_KEY)
    streamlit_js_eval(
        js_expressions=f"localStorage.removeItem({storage_key}); 'cleared';",
        key=f"auth_local_storage_clear_{int(time.time() * 1000)}",
    )
    record_auth_debug(local_storage_exists=False)


def fetch_supabase_user(access_token: str) -> dict | None:
    supabase_url, _ = get_supabase_auth_config()
    response = requests.get(
        f"{supabase_url}/auth/v1/user",
        headers=get_auth_headers(access_token),
        timeout=20,
    )
    if response.status_code >= 400:
        return None
    return response.json()


def refresh_auth_session(refresh_token: str) -> dict | None:
    record_auth_debug(refresh_attempted=True)
    data = supabase_auth_request(
        "token?grant_type=refresh_token",
        {"refresh_token": refresh_token},
    )
    user = save_auth_session(data)
    record_auth_debug(refresh_success=bool(user))
    return user


def refresh_supabase_session_if_needed(payload: dict) -> dict | None:
    access_token = payload.get("access_token")
    refresh_token = payload.get("refresh_token")
    expires_at = int(payload.get("expires_at") or 0)

    if access_token and (not expires_at or expires_at > int(time.time()) + 60):
        user_data = fetch_supabase_user(access_token)
        if user_data:
            payload["user"] = normalize_supabase_user(user_data)
            return apply_auth_payload(payload, restored=True)

    if refresh_token:
        try:
            return refresh_auth_session(refresh_token)
        except Exception:
            clear_auth_session()
            st.session_state["dentpilot_session_expired_message"] = "\u767b\u5f55\u5df2\u8fc7\u671f\uff0c\u8bf7\u91cd\u65b0\u767b\u5f55\u3002"
            return None

    return None


def restore_supabase_session_from_tokens(payload: dict | None) -> dict | None:
    if not payload:
        return None
    record_auth_debug(restore_attempted=True)
    return refresh_supabase_session_if_needed(payload)


def clear_auth_session() -> None:
    keys_to_clear = (
        *AUTH_SESSION_KEYS,
        "auth_local_storage_pending",
        "auth_local_storage_attempts",
        "auth_local_storage_recently_saved",
        "dentpilot_auth_restored_message",
        "dentpilot_session_expired_message",
    )
    for key in keys_to_clear:
        st.session_state.pop(key, None)
    clear_auth_session_from_local_storage()
    record_auth_debug(session_user=False, current_user_email=None)


def load_auth_session() -> dict | None:
    user = st.session_state.get("auth_user") or st.session_state.get("dentpilot_user")
    if user:
        record_auth_debug(session_user=True, current_user_email=user.get("email"))
        return user

    payload = load_auth_session_from_local_storage()
    return restore_supabase_session_from_tokens(payload)


def get_current_user() -> dict | None:
    if demo_login_enabled():
        user = get_demo_user()
        st.session_state["auth_user"] = user
        st.session_state["dentpilot_user"] = user
        return user
    return load_auth_session()


def sign_out_supabase() -> None:
    access_token = st.session_state.get("dentpilot_access_token") or (st.session_state.get("auth_session") or {}).get("access_token")
    if not access_token:
        return
    supabase_url, _ = get_supabase_auth_config()
    try:
        requests.post(
            f"{supabase_url}/auth/v1/logout",
            headers=get_auth_headers(access_token),
            timeout=10,
        )
    except Exception:
        pass


def render_auth_gate() -> None:
    user = get_current_user()
    if user:
        restored_message = st.session_state.pop("dentpilot_auth_restored_message", None)
        if restored_message:
            st.success(restored_message)
        return

    if st.session_state.pop("auth_local_storage_pending", False):
        st.info(t("restoring_login"))
        render_auth_debug()
        if int(st.session_state.get("auth_local_storage_attempts", 0)) <= 2:
            time.sleep(0.25)
            st.rerun()
        if not st.button(t("continue_login"), use_container_width=True):
            st.caption(
                "If this screen does not move forward, click Continue Login or refresh once."
                if get_ui_lang() == "en"
                else "如果页面没有自动继续，请点击继续登录或刷新一次。"
            )
            st.stop()

    st.markdown(
        f"""
        <section class="hero auth-hero">
            <div class="eyebrow">DentPilot AI &#36134;&#21495;&#31995;&#32479;</div>
            <h1 class="hero-title">{t("login_title")}</h1>
            <p class="hero-copy">{t("login_intro")}</p>
        </section>
        """,
        unsafe_allow_html=True,
    )

    expired_message = st.session_state.pop("dentpilot_session_expired_message", None)
    if expired_message:
        st.warning(expired_message)

    render_auth_debug()

    login_tab, register_tab = st.tabs([t("login"), t("register")])
    with login_tab:
        with st.form("dentpilot_login_form"):
            email = st.text_input(t("email"), key="login_email")
            password = st.text_input(t("password"), type="password", key="login_password")
            submitted = st.form_submit_button(t("login"), use_container_width=True)
        if submitted:
            try:
                data = sign_in_with_email(email.strip(), password)
                user = save_auth_session(data)
                if not user:
                    raise RuntimeError("Supabase \u672a\u8fd4\u56de\u53ef\u4fdd\u5b58\u7684\u767b\u5f55\u4f1a\u8bdd\u3002")
                st.success("Login successful." if get_ui_lang() == "en" else "\u767b\u5f55\u6210\u529f")
                st.rerun()
            except Exception as exc:
                st.error(f"\u767b\u5f55\u5931\u8d25\uff1a{exc}")

    with register_tab:
        with st.form("dentpilot_register_form"):
            email = st.text_input(t("email"), key="register_email")
            password = st.text_input(t("password"), type="password", key="register_password")
            submitted = st.form_submit_button(t("register"), use_container_width=True)
        if submitted:
            try:
                data = sign_up_with_email(email.strip(), password)
                user = save_auth_session(data)
                if user:
                    st.success("Registration successful. You are now logged in." if get_ui_lang() == "en" else "\u6ce8\u518c\u6210\u529f\uff0c\u5df2\u81ea\u52a8\u767b\u5f55\u3002")
                    st.rerun()
                st.success(
                    "Registration successful. If email confirmation is enabled, check your email first, then log in."
                    if get_ui_lang() == "en"
                    else "\u6ce8\u518c\u6210\u529f\u3002\u5982\u679c Supabase \u5f00\u542f\u90ae\u7bb1\u786e\u8ba4\uff0c\u8bf7\u5148\u67e5\u6536\u90ae\u4ef6\uff0c\u7136\u540e\u518d\u767b\u5f55\u3002"
                )
            except Exception as exc:
                st.error(f"\u6ce8\u518c\u5931\u8d25\uff1a{exc}")

    st.stop()


def render_sidebar_account() -> None:
    user = get_current_user() or {}
    email = user.get("email") or "\u5f53\u524d\u7528\u6237"
    st.markdown(f"### {t('current_user')}")
    st.caption(str(email))
    if st.button(t("logout"), use_container_width=True):
        sign_out_supabase()
        clear_selected_mode()
        clear_auth_session()
        st.success("\u5df2\u9000\u51fa\u767b\u5f55\u3002")
        st.stop()
    st.markdown("---")


def render_my_learning_summary() -> None:
    try:
        summary = get_user_learning_summary()
    except Exception as exc:
        st.error(f"无法读取学习记录：{exc}")
        return

    usage = summary.get("usage") or {}
    study_records = summary.get("study_pack_records") or []
    written_records = summary.get("written_exam_attempts") or []
    clinical_records = summary.get("clinical_case_attempts") or []
    oral_records = summary.get("oral_exam_attempts") or []
    weaknesses = summary.get("weaknesses") or []

    st.markdown("### 我的学习记录")
    st.caption(f"今日口语用时：{float(usage.get('voice_minutes_used') or 0):.1f} 分钟")
    average_score = summary.get("average_score")
    st.caption(f"近期平均分：{average_score if average_score is not None else '暂无'}")
    st.caption(
        f"记录数：复习包 {len(study_records)} / 笔试 {len(written_records)} / 病例 {len(clinical_records)} / 口试 {len(oral_records)}"
    )
    if written_records:
        st.caption(f"最近笔试：{written_records[0].get('topic') or written_records[0].get('subject')}")
    if clinical_records:
        st.caption(f"最近病例：{clinical_records[0].get('case_title') or 'Clinical case'}")
    if oral_records:
        st.caption(f"最近口试：{oral_records[0].get('topic') or oral_records[0].get('subject')}")
    if weaknesses:
        weak_topics = "，".join(str(item.get("topic", "")) for item in weaknesses[:3] if item.get("topic"))
        st.caption(f"薄弱题目：{weak_topics or '暂无'}")
    if not any([study_records, written_records, clinical_records, oral_records]):
        st.caption("还没有学习记录，先做一组训练吧。")
    st.markdown("---")


def admin_email_allowlist() -> set[str]:
    raw = read_config_value("ADMIN_EMAILS", read_config_value("NEXT_PUBLIC_ADMIN_EMAILS", ""))
    return {email.strip().lower() for email in raw.split(",") if email.strip()}


def is_current_user_admin() -> bool:
    user = get_current_user() or {}
    email = str(user.get("email") or "").strip().lower()
    return bool(email and email in admin_email_allowlist())


def get_service_role_key() -> str:
    return read_config_value("SUPABASE_SERVICE_ROLE_KEY", "")


def supabase_admin_request(table: str, query: str = "select=*") -> list[dict]:
    service_role_key = get_service_role_key()
    if not service_role_key:
        raise RuntimeError("SUPABASE_SERVICE_ROLE_KEY未配置，请先在服务器端配置服务端密钥。")
    supabase_url, _ = get_supabase_auth_config()
    response = requests.get(
        f"{supabase_url}/rest/v1/{table}?{query}",
        headers={
            "apikey": service_role_key,
            "Authorization": f"Bearer {service_role_key}",
            "Accept": "application/json",
        },
        timeout=30,
    )
    if response.status_code >= 400:
        raise RuntimeError(response.text)
    return response.json() if response.content else []


def render_admin_dashboard() -> None:
    if not is_current_user_admin():
        return
    with st.expander("管理员总览", expanded=False):
        if not get_service_role_key():
            st.warning("管理员统计需要在服务器环境变量或 Streamlit Secrets 中配置 SUPABASE_SERVICE_ROLE_KEY。")
            return
        try:
            profiles = supabase_admin_request("profiles", "select=id,email,created_at&limit=1000")
            study = supabase_admin_request("study_pack_records", "select=id,user_id,created_at&limit=1000")
            written = supabase_admin_request("written_exam_attempts", "select=id,user_id,score,created_at&limit=1000")
            clinical = supabase_admin_request("clinical_case_attempts", "select=id,user_id,score,created_at&limit=1000")
            oral = supabase_admin_request("oral_exam_attempts", "select=id,user_id,score,created_at&limit=1000")
            usage = supabase_admin_request("usage_limits", f"select=*&date=eq.{time.strftime('%Y-%m-%d')}&limit=1000")
        except Exception as exc:
            st.error(f"管理员数据读取失败：{exc}")
            return

        col1, col2, col3 = st.columns(3)
        col1.metric("Users", len(profiles))
        col2.metric("Study Packs", len(study))
        col3.metric("Voice minutes today", round(sum(float(row.get("voice_minutes_used") or 0) for row in usage), 1))
        col4, col5, col6 = st.columns(3)
        col4.metric("Written Exams", len(written))
        col5.metric("Clinical Cases", len(clinical))
        col6.metric("Oral Exams", len(oral))
        if usage:
            st.dataframe(pd.DataFrame(usage), use_container_width=True, hide_index=True)


ALLOWED_MODE_KEYS = {
    "study_pack": "Study Pack",
    "written_exam": "AI Written Exam",
    "clinical_case": "Clinical Case",
    "weakness_analysis": "Weakness Analysis",
    "realtime_oral_exam": "Realtime Oral Exam",
}
MODE_LABEL_TO_KEY = {label: key for key, label in ALLOWED_MODE_KEYS.items()}
LOCAL_STORAGE_MODE_KEY = "DENTPILOT_SELECTED_MODE"


def make_local_storage_mode_set_js(mode_key: str) -> str:
    storage_key = json.dumps(LOCAL_STORAGE_MODE_KEY)
    storage_value = json.dumps(mode_key)
    return f"localStorage.setItem({storage_key}, {storage_value}); 'saved';"


def save_selected_mode(mode_key: str) -> None:
    if mode_key not in ALLOWED_MODE_KEYS:
        mode_key = "study_pack"
    st.session_state["selected_mode"] = mode_key
    streamlit_js_eval(
        js_expressions=make_local_storage_mode_set_js(mode_key),
        key=f"selected_mode_save_{mode_key}",
    )


def load_selected_mode() -> str:
    session_mode = st.session_state.get("selected_mode")
    if session_mode in ALLOWED_MODE_KEYS:
        return session_mode

    raw_value = streamlit_js_eval(
        js_expressions=(
            f"localStorage.getItem({json.dumps(LOCAL_STORAGE_MODE_KEY)}) || 'study_pack';"
        ),
        key="selected_mode_load",
    )
    if raw_value in ALLOWED_MODE_KEYS:
        st.session_state["selected_mode"] = raw_value
        return raw_value
    st.session_state["selected_mode"] = "study_pack"
    return "study_pack"


def clear_selected_mode() -> None:
    st.session_state.pop("selected_mode", None)
    streamlit_js_eval(
        js_expressions=f"localStorage.removeItem({json.dumps(LOCAL_STORAGE_MODE_KEY)}); 'cleared';",
        key=f"selected_mode_clear_{int(time.time() * 1000)}",
    )


def current_user() -> dict:
    user = get_current_user() or {}
    if not user.get("id"):
        raise RuntimeError("请先登录。")
    return user


def current_user_id() -> str:
    return str(current_user()["id"])


def current_access_token() -> str:
    token = st.session_state.get("dentpilot_access_token") or (st.session_state.get("auth_session") or {}).get("access_token")
    if not token:
        raise RuntimeError("登录会话已过期，请重新登录。")
    return str(token)


def supabase_rest_url(table: str, query: str = "") -> str:
    supabase_url, _ = get_supabase_auth_config()
    base = f"{supabase_url}/rest/v1/{table}"
    return f"{base}?{query}" if query else base


def supabase_rest_request(method: str, table: str, query: str = "", payload=None):
    token = current_access_token()
    headers = get_auth_headers(token)
    headers["Accept"] = "application/json"
    if method.upper() in {"POST", "PATCH"}:
        headers["Prefer"] = "return=representation"

    response = requests.request(
        method,
        supabase_rest_url(table, query),
        headers=headers,
        json=payload,
        timeout=30,
    )
    if response.status_code >= 400:
        try:
            data = response.json()
        except Exception:
            data = {}
        message = data.get("message") or data.get("details") or data.get("hint") or response.text
        raise RuntimeError(message)
    if not response.content:
        return None
    return response.json()


def select_own_rows(table: str, limit: int = 20, extra_query: str = "", order_by: str = "created_at.desc") -> list[dict]:
    user_id = quote(current_user_id(), safe="")
    query = f"select=*&user_id=eq.{user_id}&order={order_by}&limit={limit}"
    if extra_query:
        query = f"{query}&{extra_query}"
    data = supabase_rest_request("GET", table, query)
    return data if isinstance(data, list) else []


def insert_own_row(table: str, payload: dict) -> dict | None:
    payload = {**payload, "user_id": current_user_id()}
    data = supabase_rest_request("POST", table, "select=*", payload)
    if isinstance(data, list) and data:
        return data[0]
    return None


def save_study_pack_record(subject: str, source_title: str, source_text: str, pack: dict, markdown_export: str) -> dict | None:
    return insert_own_row(
        "study_pack_records",
        {
            "subject": subject,
            "source_title": source_title,
            "source_text": source_text,
            "generated_pack": pack,
            "markdown_export": markdown_export,
        },
    )


def get_recent_study_pack_records(limit: int = 10) -> list[dict]:
    return select_own_rows("study_pack_records", limit)


def save_written_exam_attempt(attempt: dict) -> dict | None:
    result = attempt.get("result") or attempt.get("grading_result") or {}
    return insert_own_row(
        "written_exam_attempts",
        {
            "session_id": attempt.get("session_id"),
            "topic": attempt.get("topic") or attempt.get("subject"),
            "subject": attempt.get("subject"),
            "difficulty": attempt.get("difficulty"),
            "course_context": attempt.get("course_context"),
            "question": attempt.get("question"),
            "question_type": attempt.get("question_type"),
            "options": attempt.get("options") or [],
            "correct_answer": attempt.get("correct_answer"),
            "student_answer": attempt.get("answer") or attempt.get("student_answer"),
            "model_answer": attempt.get("model_answer"),
            "score": safe_float(result.get("score")),
            "feedback": result.get("chinese_feedback") or result.get("feedback"),
            "covered_points": result.get("covered_points") or [],
            "missing_points": result.get("missing_points") or [],
        },
    )


def get_recent_written_exam_attempts(limit: int = 20) -> list[dict]:
    return select_own_rows("written_exam_attempts", limit)


def save_clinical_case_attempt(attempt: dict) -> dict | None:
    result = attempt.get("result") or {}
    return insert_own_row(
        "clinical_case_attempts",
        {
            "case_title": attempt.get("case_title"),
            "case_data": attempt.get("case_data") or {},
            "student_answer": attempt.get("answer") or attempt.get("student_answer"),
            "score": safe_float(result.get("score")),
            "diagnosis_score": safe_float(result.get("diagnosis_score")),
            "treatment_score": safe_float(result.get("treatment_score")),
            "missing_points": result.get("missing_points") or [],
            "feedback": result.get("chinese_feedback") or result.get("feedback"),
        },
    )


def get_recent_clinical_case_attempts(limit: int = 20) -> list[dict]:
    return select_own_rows("clinical_case_attempts", limit)


def save_oral_exam_attempt(attempt: dict) -> dict | None:
    result = attempt.get("result") or {}
    return insert_own_row(
        "oral_exam_attempts",
        {
            "session_id": attempt.get("session_id"),
            "question": attempt.get("question"),
            "student_answer": attempt.get("answer") or attempt.get("student_answer"),
            "score": safe_float(result.get("score")),
            "covered_points": result.get("covered_points") or [],
            "missing_points": result.get("missing_points") or [],
            "feedback": result.get("chinese_feedback") or result.get("feedback"),
            "topic": attempt.get("topic"),
            "subject": attempt.get("subject"),
            "difficulty": attempt.get("difficulty"),
        },
    )


def get_recent_oral_exam_attempts(limit: int = 20) -> list[dict]:
    return select_own_rows("oral_exam_attempts", limit)


def get_usage_today() -> dict | None:
    user_id = quote(current_user_id(), safe="")
    today = time.strftime("%Y-%m-%d")
    query = f"select=*&user_id=eq.{user_id}&date=eq.{today}&limit=1"
    data = supabase_rest_request("GET", "usage_limits", query)
    if isinstance(data, list) and data:
        return data[0]
    inserted = insert_own_row("usage_limits", {"date": today})
    return inserted


def increment_usage(kind: str, amount: float = 1) -> dict | None:
    usage = get_usage_today()
    if not usage:
        return None
    field_map = {
        "voice_minutes": "voice_minutes_used",
        "text_exam": "text_exam_count",
        "study_pack": "study_pack_count",
        "case": "case_count",
    }
    field = field_map.get(kind)
    if not field:
        raise ValueError(f"Unsupported usage kind: {kind}")
    if field == "voice_minutes_used":
        next_value = safe_float(usage.get(field), 0.0) + safe_float(amount, 0.0)
    else:
        next_value = safe_int(usage.get(field), 0) + safe_int(amount, 0)
    row_id = quote(str(usage["id"]), safe="")
    user_id = quote(current_user_id(), safe="")
    today = time.strftime("%Y-%m-%d")
    data = supabase_rest_request(
        "PATCH",
        "usage_limits",
        f"id=eq.{row_id}&user_id=eq.{user_id}&date=eq.{today}&select=*",
        {field: next_value},
    )
    return data[0] if isinstance(data, list) and data else None


def get_user_weakness_summary(limit: int = 20) -> list[dict]:
    return select_own_rows("user_weaknesses", limit, order_by="last_seen_at.desc")


def update_user_weaknesses_from_attempt(subject: str, topic: str, missing_points: list, score: float | None = None) -> None:
    for point in (missing_points or [topic])[:5]:
        insert_own_row(
            "user_weaknesses",
            {
                "subject": subject,
                "topic": str(point or topic),
                "weakness_type": "missing_point" if missing_points else "practice_topic",
                "score_avg": safe_float(score),
                "attempt_count": safe_int(1),
            },
        )


def get_user_learning_summary() -> dict:
    written = get_recent_written_exam_attempts(50)
    clinical = get_recent_clinical_case_attempts(50)
    oral = get_recent_oral_exam_attempts(50)
    study = get_recent_study_pack_records(50)
    weaknesses = get_user_weakness_summary(20)
    usage = get_usage_today()
    scores = [
        float(item.get("score"))
        for item in [*written, *clinical, *oral]
        if item.get("score") is not None
    ]
    return {
        "study_pack_records": study,
        "written_exam_attempts": written,
        "clinical_case_attempts": clinical,
        "oral_exam_attempts": oral,
        "weaknesses": weaknesses,
        "usage": usage,
        "average_score": round(sum(scores) / len(scores), 1) if scores else None,
    }


def load_persistent_learning_records() -> None:
    try:
        st.session_state["study_pack_records"] = get_recent_study_pack_records()
    except Exception as exc:
        st.session_state["study_pack_records_error"] = str(exc)
    try:
        st.session_state["oral_exam_history"] = get_recent_written_exam_attempts()
    except Exception as exc:
        st.session_state["oral_exam_history_error"] = str(exc)
    try:
        st.session_state["clinical_case_history"] = get_recent_clinical_case_attempts()
    except Exception as exc:
        st.session_state["clinical_case_history_error"] = str(exc)
    try:
        st.session_state["realtime_oral_history"] = get_recent_oral_exam_attempts()
    except Exception as exc:
        st.session_state["realtime_oral_history_error"] = str(exc)


def extract_pdf_text(uploaded_file) -> tuple[str, dict]:
    uploaded_file.seek(0)
    pdf = PdfReader(uploaded_file)
    pages: list[str] = []
    failed_pages: list[int] = []
    empty_pages: list[int] = []
    for page_index, page in enumerate(pdf.pages, start=1):
        try:
            page_text = page.extract_text() or ""
        except Exception:
            page_text = ""
            failed_pages.append(page_index)
        if page_text.strip():
            pages.append(f"\n\n--- Page {page_index} ---\n{page_text.strip()}")
        elif page_index not in failed_pages:
            empty_pages.append(page_index)
    text = "\n".join(pages).strip()
    report = {
        "filename": getattr(uploaded_file, "name", "uploaded.pdf"),
        "total_pages": len(pdf.pages),
        "successful_pages": len(pdf.pages) - len(failed_pages) - len(empty_pages),
        "failed_pages": failed_pages,
        "empty_pages": empty_pages,
        "char_count": len(text),
        "preview": text[:1000],
    }
    return text, report


def extract_docx_text(uploaded_file) -> tuple[str, int]:
    uploaded_file.seek(0)
    document = Document(uploaded_file)
    blocks = []
    for paragraph in document.paragraphs:
        paragraph_text = paragraph.text.strip()
        if paragraph_text:
            blocks.append(paragraph_text)
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                blocks.append(" | ".join(cells))
    return "\n\n".join(blocks), len(document.paragraphs)


def extract_pptx_text(uploaded_file) -> tuple[str, dict]:
    if Presentation is None:
        raise RuntimeError("当前环境缺少 python-pptx，无法解析 PowerPoint 文件。")
    uploaded_file.seek(0)
    presentation = Presentation(uploaded_file)
    blocks: list[str] = []
    slides_with_text = 0

    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_blocks: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text_frame:
                text = "\n".join(
                    paragraph.text.strip()
                    for paragraph in shape.text_frame.paragraphs
                    if paragraph.text.strip()
                )
                if text:
                    slide_blocks.append(text)
            if getattr(shape, "has_table", False):
                rows = []
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        rows.append(" | ".join(cells))
                if rows:
                    slide_blocks.append("\n".join(rows))
        if slide_blocks:
            slides_with_text += 1
            blocks.append(f"\n\n--- Slide {slide_index} ---\n" + "\n\n".join(slide_blocks))

    text = "\n".join(blocks).strip()
    report = {
        "filename": getattr(uploaded_file, "name", "uploaded.pptx"),
        "total_slides": len(presentation.slides),
        "successful_slides": slides_with_text,
        "char_count": len(text),
        "preview": text[:1000],
    }
    return text, report


def extract_txt_text(uploaded_file) -> tuple[str, int]:
    uploaded_file.seek(0)
    raw_bytes = uploaded_file.read()
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "latin-1"):
        try:
            text = raw_bytes.decode(encoding)
            break
        except UnicodeDecodeError:
            continue
    else:
        text = raw_bytes.decode("utf-8", errors="ignore")
    lines = [line for line in text.splitlines() if line.strip()]
    return text.strip(), len(lines)


def get_pdf_font_name() -> str:
    font_candidates = [
        Path("C:/Windows/Fonts/msyh.ttc"),
        Path("C:/Windows/Fonts/simhei.ttf"),
        Path("C:/Windows/Fonts/simsun.ttc"),
        Path("C:/Windows/Fonts/arialuni.ttf"),
    ]
    for font_path in font_candidates:
        if font_path.exists():
            font_name = "DentPilotCJK"
            if font_name not in pdfmetrics.getRegisteredFontNames():
                pdfmetrics.registerFont(TTFont(font_name, str(font_path)))
            return font_name
    font_name = "STSong-Light"
    if font_name not in pdfmetrics.getRegisteredFontNames():
        pdfmetrics.registerFont(UnicodeCIDFont(font_name))
    return font_name


def as_paragraph_text(value) -> str:
    text_value = str(value or "")
    return (
        text_value.replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace("\n", "<br/>")
    )


def build_study_pack_pdf(pack: dict, subject: str) -> bytes:
    """Build a simple PDF report for study pack outputs."""
    buffer = io.BytesIO()
    font_name = get_pdf_font_name()
    styles = getSampleStyleSheet()

    body = ParagraphStyle(
        "DentPilotBody",
        parent=styles["BodyText"],
        fontName=font_name,
        fontSize=10,
        leading=15,
        spaceAfter=8,
    )
    title_style = ParagraphStyle(
        "DentPilotTitle",
        parent=styles["Title"],
        fontName=font_name,
        fontSize=22,
        leading=28,
        textColor=colors.HexColor("#0f172a"),
        spaceAfter=10,
    )
    section = ParagraphStyle(
        "DentPilotSection",
        parent=styles["Heading2"],
        fontName=font_name,
        fontSize=14,
        leading=18,
        textColor=colors.HexColor("#1e40af"),
        spaceBefore=14,
        spaceAfter=8,
    )

    doc = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        rightMargin=0.55 * inch,
        leftMargin=0.55 * inch,
        topMargin=0.6 * inch,
        bottomMargin=0.6 * inch,
        title="DentPilot AI Study Pack",
    )

    detected_sections = pack.get("coverage_report", {}).get("detected_sections", 0)
    generated_sections = pack.get("coverage_report", {}).get("generated_sections", 0)
    coverage_percent = pack.get("coverage_report", {}).get("coverage_percent", 100)

    story = [
        Paragraph("DentPilot AI Study Pack", title_style),
        Paragraph(f"Subject: {as_paragraph_text(subject)}", body),
        Paragraph(
            f"Generation mode: {as_paragraph_text(pack.get('generation_depth', 'standard'))}",
            body,
        ),
        Spacer(1, 8),
    ]

    modules = pack.get("study_modules", [])
    processing_details = pack.get("processing_details", {})
    if modules:
        story.extend(
            [
                Paragraph("Coverage Report", section),
                Paragraph(
                    f"Extracted characters: {processing_details.get('extracted_char_count', 0)}<br/>"
                    f"Total chunks: {coverage.get('total_chunks', detected_sections)}<br/>"
                    f"Processed chunks: {processing_details.get('processed_chunk_count', generated_sections)}<br/>"
                    f"Failed chunks: {processing_details.get('failed_chunk_count', 0)}<br/>"
                    f"Model: {as_paragraph_text(processing_details.get('model_used', ''))}<br/>"
                    f"Early topics: {as_paragraph_text(', '.join(processing_details.get('early_topics', [])))}<br/>"
                    f"Middle topics: {as_paragraph_text(', '.join(processing_details.get('middle_topics', [])))}<br/>"
                    f"Late topics: {as_paragraph_text(', '.join(processing_details.get('late_topics', [])))}",
                    body,
                ),
                Paragraph(
                    f"Detected {detected_sections} sections, generated {generated_sections} modules."
                    f" Coverage {coverage_percent}%.",
                    body,
                ),
                Spacer(1, 6),
            ]
        )

        for module in modules:
            title = module.get("title") or module.get("question_title") or "Study module"
            section_title = module.get("section_number")
            if section_title:
                title = f"Question {section_title}: {title}"
            story.append(Paragraph(as_paragraph_text(title), section_style=section if False else section))
            story.append(Paragraph("<b>Chinese core explanation</b>", body))
            story.append(Paragraph(as_paragraph_text(module.get("chinese_core_explanation", "")), body))

            story.append(Paragraph("<b>Must-know points</b>", body))
            for item in module.get("must_know", []):
                story.append(Paragraph(as_paragraph_text(f"- {item}"), body))

            story.append(Paragraph("<b>Common mistakes</b>", body))
            for item in module.get("common_mistakes", []):
                story.append(Paragraph(as_paragraph_text(f"- {item}"), body))

            story.append(Paragraph("<b>Likely oral exam questions</b>", body))
            for item in module.get("oral_exam_questions", []):
                story.append(Paragraph(as_paragraph_text(f"- {item}"), body))

            story.append(Paragraph("<b>Short answer template</b>", body))
            story.append(Paragraph(as_paragraph_text(module.get("short_answer_template", "")), body))

            story.append(Paragraph("<b>Follow-up questions</b>", body))
            for item in module.get("follow_up_questions", []):
                story.append(Paragraph(as_paragraph_text(f"- {item}"), body))

            story.append(Paragraph("<b>Anki cards</b>", body))
            for card in module.get("flashcards", []):
                story.append(
                    Paragraph(
                        as_paragraph_text(
                            f"{card.get('type', 'card')}: {card.get('front', '')} -> {card.get('back', '')}"
                        ),
                        body,
                    )
                )
            story.append(Spacer(1, 6))

        story.append(Paragraph("Quiz", section))
        for index, q in enumerate(pack.get("quiz", []), start=1):
            options = q.get("options", [])
            option_text = "<br/>".join(as_paragraph_text(option) for option in options)
            question_text = (
                f"<b>Q{index}. [{as_paragraph_text(q.get('question_type', 'quiz'))}] "
                f"{as_paragraph_text(q.get('question', ''))}</b><br/>"
                f"{option_text}<br/>"
                f"<b>Answer:</b> {as_paragraph_text(q.get('answer', ''))}<br/>"
                f"<b>Explanation:</b> {as_paragraph_text(q.get('explanation_zh', ''))}"
            )
            story.append(Paragraph(question_text, body))
    else:
        story.append(Paragraph("Study Notes", section))
        story.append(Paragraph(as_paragraph_text(pack.get("chinese_explanation", "")), body))

    doc.build(story)
    return buffer.getvalue()

def build_study_pack_markdown(pack: dict, subject: str) -> bytes:
    lines = [
        "# DentPilot AI Study Pack",
        "",
        f"- Subject: {subject}",
        f"- Generation depth: {pack.get('generation_depth', 'standard')}",
        "",
    ]

    modules = pack.get("study_modules", [])
    coverage = pack.get("coverage_report", {})
    processing_details = pack.get("processing_details", {})

    if modules:
        lines.extend(
            [
                "## Coverage Report",
                "",
                f"- Extracted characters: {processing_details.get('extracted_char_count', 0)}",
                f"- Total chunks: {coverage.get('total_chunks', len(modules))}",
                f"- Processed chunks: {processing_details.get('processed_chunk_count', len(modules))}",
                f"- Failed chunks: {processing_details.get('failed_chunk_count', 0)}",
                f"- Model: {processing_details.get('model_used', '')}",
                f"- Coverage status: {'complete' if not coverage.get('has_missing') else 'partial'}",
                f"- Early-section topics: {', '.join(processing_details.get('early_topics', []))}",
                f"- Middle-section topics: {', '.join(processing_details.get('middle_topics', []))}",
                f"- Late-section topics: {', '.join(processing_details.get('late_topics', []))}",
                "",
                f"- Detected sections: {coverage.get('detected_sections', len(modules))}",
                f"- Generated sections: {coverage.get('generated_sections', len(modules))}",
                f"- Coverage: {coverage.get('coverage_percent', 100)}%",
                "",
            ]
        )
        for module in modules:
            lines.extend(
                [
                    f"### Question {module.get('section_number')}: {module.get('title', '')}",
                    "",
                    "#### 中文核心讲解",
                    str(module.get("chinese_core_explanation", "")),
                    "",
                    "#### Must-know points",
                ]
            )
            lines.extend(f"- {item}" for item in module.get("must_know", []))
            lines.extend(["", "#### Common mistakes"])
            lines.extend(f"- {item}" for item in module.get("common_mistakes", []))
            lines.extend(["", "#### Likely oral exam questions"])
            lines.extend(f"- {item}" for item in module.get("oral_exam_questions", []))
            lines.extend([
                "",
                "#### Short answer template",
                str(module.get("short_answer_template", "")),
                "",
                "#### Follow-up questions",
            ])
            lines.extend(f"- {item}" for item in module.get("follow_up_questions", []))

            lines.append("")
            lines.extend(["#### Quiz / Questions:"])
            for q in module.get("quiz", []):
                lines.append(f"- [{q.get('question_type', 'quiz')}] {q.get('question', '')}")
                if q.get("answer"):
                    lines.append(f"  Answer: {q.get('answer')}")
                if q.get("explanation_zh"):
                    lines.append(f"  Explanation: {q.get('explanation_zh')}")
            lines.append("")

            lines.append("#### Anki cards")
            for card in module.get("flashcards", []):
                lines.append(f"- [{card.get('type', 'concept')}] {card.get('front', '')} -> {card.get('back', '')}")
            lines.append("")
    else:
        lines.extend(["## Study Notes", "", str(pack.get("chinese_explanation", "")), ""])

    glossary = pack.get("glossary", [])
    if glossary:
        lines.extend(["## Glossary", ""])
        for term in glossary:
            lines.append(
                f"- {term.get('english', '')} / {term.get('chinese', '')}: {term.get('definition', '')}"
            )

    if pack.get("quiz"):
        lines.extend(["", "## Quiz", ""])
        for index, q in enumerate(pack.get("quiz", []), start=1):
            lines.extend(
                [
                    f"### Q{index}. {q.get('question', '')}",
                    "",
                    *[f"- {option}" for option in q.get("options", [])],
                    "",
                    f"Answer: {q.get('answer', '')}",
                    f"Explanation: {q.get('explanation_zh', '')}",
                    "",
                ]
            )

    lines.extend(["## Anki Cards", ""])
    for card in pack.get("flashcards", []):
        lines.extend(
            [
                f"- Front: {card.get('front', '')}",
                f"  Back: {card.get('back', '')}",
                f"  Type: {card.get('type', 'concept')}",
                "",
            ]
        )

    lines.extend(["## 考前总结", "", str(pack.get("exam_summary", "")), ""])
    return "\n".join(lines).encode("utf-8")

def build_anki_csv_bytes(pack: dict, subject: str) -> bytes:
    output = io.StringIO()
    writer = csv.writer(output)
    writer.writerow(["Front", "Back", "Tags"])
    for card in pack.get("flashcards", []):
        writer.writerow([
            card.get("front", ""),
            card.get("back", ""),
            f"medstudy::{subject}::{card.get('type', 'concept')}",
        ])
    return output.getvalue().encode("utf-8-sig")


def render_dental_loader(message: str, submessage: str | None = None, *, container: st.delta_generator.DeltaGenerator | None = None) -> None:
    if not st.session_state.get("_dentpilot_dental_loader_css_injected"):
        st.session_state["_dentpilot_dental_loader_css_injected"] = True
        st.markdown(
            """
            <style>
            .dentpilot-loader-container {
                display: flex;
                flex-direction: column;
                align-items: center;
                gap: 12px;
                margin: 12px 0 8px;
                color: #0f172a;
            }
            .dentpilot-loader-stage {
                position: relative;
                width: 104px;
                height: 104px;
                display: flex;
                align-items: center;
                justify-content: center;
                border-radius: 999px;
            }
            .dentpilot-pulse-ring {
                position: absolute;
                width: 102px;
                height: 102px;
                border-radius: 50%;
                border: 2px solid rgba(20, 184, 166, 0.45);
                animation: dentpilot-ring 2.2s ease-in-out infinite;
            }
            .dentpilot-tooth-wrap {
                position: absolute;
                width: 72px;
                height: 72px;
                display: flex;
                align-items: center;
                justify-content: center;
                border-radius: 50%;
                animation: dentpilot-tooth-float 2.4s ease-in-out infinite;
                filter: drop-shadow(0 8px 22px rgba(30, 64, 175, 0.24));
            }
            .dentpilot-tooth-icon {
                width: 48px;
                height: 54px;
            }
            .dentpilot-tooth-path {
                fill: #0f172a;
                transform-origin: center;
                animation: dentpilot-tooth-tilt 2.8s ease-in-out infinite;
            }
            .dentpilot-heartline {
                width: 90px;
                height: 20px;
                position: relative;
            }
            .dentpilot-heartline svg {
                width: 100%;
                height: 100%;
            }
            .dentpilot-heartline-path {
                fill: none;
                stroke: #0ea5e9;
                stroke-width: 2.6;
                stroke-linecap: round;
                stroke-linejoin: round;
                stroke-dasharray: 140;
                stroke-dashoffset: 140;
                animation: dentpilot-heartline 1.5s ease-in-out infinite;
            }
            .dentpilot-loader-text {
                text-align: center;
                font-weight: 600;
                color: #0f172a;
                letter-spacing: 0.02em;
            }
            .dentpilot-loader-subtext {
                margin-top: 2px;
                font-size: 0.92rem;
                color: #334155;
            }
            @keyframes dentpilot-ring {
                0% {
                    transform: scale(0.92);
                    opacity: 0.32;
                    border-color: rgba(20, 184, 166, 0.25);
                }
                50% {
                    transform: scale(1.06);
                    opacity: 0.85;
                    border-color: rgba(37, 99, 235, 0.58);
                }
                100% {
                    transform: scale(0.92);
                    opacity: 0.32;
                    border-color: rgba(20, 184, 166, 0.25);
                }
            }
            @keyframes dentpilot-tooth-float {
                0% { transform: translateY(0px) rotate(0deg); }
                50% { transform: translateY(-4px) rotate(-2deg); }
                100% { transform: translateY(0px) rotate(0deg); }
            }
            @keyframes dentpilot-tooth-tilt {
                0% { transform: rotate(0deg) scale(1); }
                50% { transform: rotate(5deg) scale(1.03); }
                100% { transform: rotate(0deg) scale(1); }
            }
            @keyframes dentpilot-heartline {
                0% { stroke-dashoffset: 140; opacity: 0.4; }
                35% { opacity: 1; }
                100% { stroke-dashoffset: 0; opacity: 0.8; }
            }
            </style>
            """,
            unsafe_allow_html=True,
        )

    safe_submessage = f"<div class='dentpilot-loader-subtext'>{submessage}</div>" if submessage else ""
    target = container or st
    target.markdown(
        f"""
        <div class="dentpilot-loader-container">
            <div class="dentpilot-loader-stage">
                <div class="dentpilot-pulse-ring"></div>
                <div class="dentpilot-tooth-wrap">
                    <svg viewBox="0 0 64 64" class="dentpilot-tooth-icon" aria-hidden="true">
                        <path
                            class="dentpilot-tooth-path"
                            d="M32 2C22.6 2 16 9.4 16 17.6c0 6.8 3.4 12.3 8.6 15.2v15.1c0 6.2-2.8 10.9-3.9 16.6-1.2 6.1 0.8 13.1 4.6 17.8 2.6 3.1 6 5.4 10 6.5 3.5 0.9 7.6 0.9 11.1 0 3.9-1.1 7.3-3.4 10-6.5 3.8-4.7 5.8-11.7 4.6-17.8-1.1-5.7-3.9-10.4-3.9-16.6V32.8C44.6 29.9 48 24.4 48 17.6 48 9.4 41.4 2 32 2Zm-4.8 55.4C21.9 57.2 18 57 18 57h28s-3.9.2-9.2 1.6c-2.3.6-5.3.6-8-.1Z"
                        />
                    </svg>
                </div>
            </div>
            <div class="dentpilot-heartline">
                <svg viewBox="0 0 90 20" preserveAspectRatio="none">
                    <polyline class="dentpilot-heartline-path" points="0,14 18,8 34,14 46,4 60,16 70,10 90,14"/>
                </svg>
            </div>
            <div class="dentpilot-loader-text">{message}{safe_submessage}</div>
        </div>
        """,
        unsafe_allow_html=True,
    )


def run_with_dental_loader(message: str, task_fn, submessage: str | None = None):
    loader_container = st.empty()
    render_dental_loader(message, submessage, container=loader_container)
    try:
        return task_fn()
    finally:
        loader_container.empty()


st.set_page_config(
    page_title="DentPilot AI",
    page_icon="🦷",
    layout="wide",
    initial_sidebar_state="expanded",
)

init_ui_language()


st.markdown(
    """
    <style>
        :root {
            --brand: #2563eb;
            --brand-2: #14b8a6;
            --brand-dark: #1e40af;
            --ink: #0f172a;
            --muted: #64748b;
            --panel: #ffffff;
            --line: #dbeafe;
            --soft: #eff6ff;
            --mint: #f0fdfa;
        }

        .stApp {
            background:
                radial-gradient(circle at 8% 0%, rgba(20, 184, 166, 0.20), transparent 28rem),
                radial-gradient(circle at 90% 12%, rgba(37, 99, 235, 0.18), transparent 32rem),
                linear-gradient(180deg, #f8fbff 0%, #eef6ff 44%, #f8fafc 100%);
            color: var(--ink);
        }

        [data-testid="stSidebar"] {
            background: linear-gradient(180deg, #0f172a 0%, #111827 62%, #0b1220 100%);
        }

        [data-testid="stSidebar"] * {
            color: #e5edf7;
        }

        [data-testid="stSidebar"] .stSelectbox label {
            color: #cbd5e1;
        }

        [data-testid="stSidebar"] hr {
            border-color: rgba(226, 232, 240, 0.18);
        }

        [data-testid="collapsedControl"] {
            top: 0.85rem;
            left: 0.85rem;
            z-index: 999999;
        }

        button[data-testid="stExpandSidebarButton"],
        [data-testid="collapsedControl"] button,
        [data-testid="stSidebarCollapseButton"] button {
            min-width: 3rem;
            min-height: 3rem;
            border: 2px solid rgba(255, 255, 255, 0.92) !important;
            border-radius: 999px !important;
            background: #2563eb !important;
            color: #ffffff !important;
            box-shadow: 0 16px 36px rgba(15, 23, 42, 0.28), 0 0 0 4px rgba(37, 99, 235, 0.18) !important;
        }

        button[data-testid="stExpandSidebarButton"]:hover,
        [data-testid="collapsedControl"] button:hover,
        [data-testid="stSidebarCollapseButton"] button:hover {
            background: #1d4ed8 !important;
            transform: translateY(-1px);
        }

        button[data-testid="stExpandSidebarButton"] span,
        button[data-testid="stExpandSidebarButton"] svg,
        [data-testid="collapsedControl"] button span,
        [data-testid="collapsedControl"] button svg {
            color: #ffffff !important;
            fill: #ffffff !important;
        }

        button[data-testid="stExpandSidebarButton"]::after,
        [data-testid="collapsedControl"] button::after {
            content: "菜单";
            position: absolute;
            left: 3.25rem;
            top: 50%;
            transform: translateY(-50%);
            padding: 0.25rem 0.5rem;
            border-radius: 999px;
            background: rgba(15, 23, 42, 0.88);
            color: #ffffff;
            font-size: 0.78rem;
            font-weight: 800;
            white-space: nowrap;
            box-shadow: 0 10px 24px rgba(15, 23, 42, 0.18);
        }

        .main .block-container {
            padding-top: 1.75rem;
            max-width: 1200px;
        }

        h1, h2, h3 {
            letter-spacing: 0;
        }

        .hero {
            position: relative;
            overflow: hidden;
            border: 1px solid rgba(37, 99, 235, 0.16);
            background:
                linear-gradient(135deg, rgba(255, 255, 255, 0.96) 0%, rgba(240, 247, 255, 0.92) 56%, rgba(240, 253, 250, 0.94) 100%);
            border-radius: 20px;
            padding: 2.25rem;
            box-shadow: 0 24px 70px rgba(15, 23, 42, 0.08);
            margin-bottom: 1.25rem;
        }

        .eyebrow {
            color: var(--brand-dark);
            font-size: 0.78rem;
            font-weight: 800;
            letter-spacing: 0.08em;
            text-transform: uppercase;
            margin-bottom: 0.6rem;
        }

        .hero-title {
            color: var(--ink);
            font-size: clamp(2.35rem, 4vw, 4.6rem);
            font-weight: 850;
            line-height: 0.98;
            margin: 0;
        }

        .hero-subtitle {
            color: var(--muted);
            font-size: 1.15rem;
            line-height: 1.65;
            max-width: 720px;
            margin-top: 1rem;
            margin-bottom: 0;
        }

        .hero-copy {
            color: #334155;
            font-size: 0.98rem;
            line-height: 1.7;
            max-width: 760px;
            margin-top: 0.75rem;
        }

        .hero-metrics {
            display: grid;
            grid-template-columns: repeat(3, minmax(0, 1fr));
            gap: 0.75rem;
            margin-top: 1.4rem;
        }

        .metric-pill {
            background: rgba(255, 255, 255, 0.72);
            border: 1px solid rgba(148, 163, 184, 0.24);
            border-radius: 12px;
            padding: 0.8rem 0.9rem;
        }

        .metric-value {
            color: var(--brand-dark);
            font-size: 1.35rem;
            font-weight: 800;
            line-height: 1;
        }

        .metric-label {
            color: var(--muted);
            font-size: 0.8rem;
            margin-top: 0.25rem;
        }

        .feature-grid {
            display: grid;
            grid-template-columns: repeat(4, minmax(0, 1fr));
            gap: 0.9rem;
            margin: 1rem 0 1.1rem;
        }

        .feature-card {
            background: var(--panel);
            border: 1px solid rgba(148, 163, 184, 0.22);
            border-radius: 12px;
            padding: 1rem;
            min-height: 132px;
            box-shadow: 0 14px 34px rgba(15, 23, 42, 0.06);
        }

        .feature-icon {
            color: var(--brand);
            font-size: 1.35rem;
            line-height: 1;
        }

        .feature-title {
            color: var(--ink);
            font-weight: 800;
            margin-top: 0.55rem;
            margin-bottom: 0.3rem;
        }

        .feature-copy {
            color: var(--muted);
            font-size: 0.9rem;
            line-height: 1.45;
        }

        .input-panel {
            background: rgba(255, 255, 255, 0.86);
            border: 1px solid rgba(148, 163, 184, 0.25);
            border-radius: 16px;
            padding: 1.25rem;
            box-shadow: 0 18px 50px rgba(15, 23, 42, 0.07);
        }

        .section-label {
            color: var(--brand-dark);
            font-weight: 800;
            font-size: 0.86rem;
            letter-spacing: 0.03em;
            text-transform: uppercase;
            margin-bottom: 0.35rem;
        }

        .section-copy {
            color: var(--muted);
            margin-bottom: 1rem;
        }

        div[data-testid="stButton"] > button {
            border-radius: 10px;
            border: 0;
            min-height: 3rem;
            font-weight: 800;
            box-shadow: 0 12px 24px rgba(37, 99, 235, 0.22);
        }

        div[data-testid="stDownloadButton"] > button {
            border-radius: 10px;
            font-weight: 700;
        }

        .result-wrap {
            background: rgba(255, 255, 255, 0.78);
            border: 1px solid rgba(148, 163, 184, 0.2);
            border-radius: 16px;
            padding: 1rem;
            margin-top: 1.4rem;
        }

        .example-grid {
            display: grid;
            grid-template-columns: minmax(0, 1.6fr) minmax(220px, 0.8fr);
            gap: 1rem;
            align-items: stretch;
            margin-bottom: 0.8rem;
        }

        .example-card,
        .workflow-card {
            background: #ffffff;
            border: 1px solid rgba(148, 163, 184, 0.24);
            border-radius: 12px;
            padding: 1rem;
        }

        .example-title {
            color: var(--ink);
            font-weight: 800;
            margin-bottom: 0.45rem;
        }

        .example-copy {
            color: var(--muted);
            font-size: 0.92rem;
            line-height: 1.55;
        }

        .footer {
            color: #64748b;
            border-top: 1px solid rgba(148, 163, 184, 0.25);
            margin-top: 2rem;
            padding: 1.2rem 0 0.25rem;
            text-align: center;
            font-size: 0.9rem;
        }

        @media (max-width: 900px) {
            .hero-metrics,
            .feature-grid,
            .example-grid {
                grid-template-columns: 1fr;
            }

            .hero {
                padding: 1.35rem;
            }

            [data-testid="collapsedControl"] {
                top: 0.75rem;
                left: 0.75rem;
            }

            button[data-testid="stExpandSidebarButton"],
            [data-testid="collapsedControl"] button {
                min-width: 3.4rem;
                min-height: 3.4rem;
            }

            button[data-testid="stExpandSidebarButton"]::after,
            [data-testid="collapsedControl"] button::after {
                left: 3.65rem;
                font-size: 0.82rem;
            }
        }
    </style>
    """,
    unsafe_allow_html=True,
)


render_auth_gate()
load_persistent_learning_records()


sample = (
    "Dental caries is a biofilm-mediated, sugar-driven, multifactorial disease that results "
    "in the demineralization of dental hard tissues. The balance between demineralization "
    "and remineralization is influenced by oral hygiene, diet, fluoride exposure, saliva, "
    "and bacterial activity. Pulpitis is inflammation of the dental pulp and may be reversible "
    "or irreversible depending on the duration and severity of symptoms."
)


def render_oral_grade_result(result: dict):
    st.markdown("### 笔试结果")
    score_col, level_col = st.columns([1, 1])
    score_col.metric("得分", f"{result.get('score', 0)}/100")
    level_col.metric("等级", result.get("level", ""))

    rubric_rows = [
        {"维度": "内容准确度", "得分": result.get("content_accuracy", 0), "权重": 30},
        {"维度": "完整性", "得分": result.get("completeness", 0), "权重": 20},
        {"维度": "临床推理", "得分": result.get("clinical_reasoning", 0), "权重": 20},
        {"维度": "语言表达", "得分": result.get("english_expression", 0), "权重": 10},
        {"维度": "考官互动", "得分": result.get("examiner_interaction", 0), "权重": 10},
        {"维度": "发音与流利度", "得分": result.get("pronunciation_fluency", 0), "权重": 10},
    ]
    st.dataframe(pd.DataFrame(rubric_rows), use_container_width=True, hide_index=True)

    col_1, col_2 = st.columns(2)
    with col_1:
        st.subheader("优点")
        strengths = result.get("strengths") or []
        if strengths:
            for item in strengths:
                st.markdown(f"- {item}")
        else:
            st.write("未识别到明显优点")
    with col_2:
        st.subheader("待改进")
        missing_points = result.get("missing_points") or []
        if missing_points:
            for item in missing_points:
                st.markdown(f"- {item}")
        else:
            st.write("未识别到明显缺失点")

    st.subheader("Corrected Answer")
    st.write(result.get("corrected_answer", ""))

    st.subheader("中文反馈")
    st.write(result.get("chinese_feedback", ""))

    st.subheader("Follow-up Question")
    st.info(result.get("follow_up_question", ""))


WRITTEN_MODES = ["日常练习", "考前模拟", "错题强化"]
WRITTEN_QUESTION_TYPES = [
    "自动混合",
    "MCQ 单选题",
    "Short Answer 简答题",
    "Case-based 病例题",
    "True / False 判断题",
    "Matching 匹配题",
]



def get_recent_wrong_written_topics(limit: int = 30) -> list[str]:
    """Return unique topics from recent written attempts with weak performance."""
    wrong_topics: list[str] = []
    seen: set[str] = set()
    for attempt in get_recent_written_exam_attempts(limit=limit * 3):
        score = safe_float(attempt.get("score"), None)
        missing_points = attempt.get("missing_points") or []
        topic = str(attempt.get("topic") or "").strip()
        if not topic:
            continue
        if (score is not None and score < 70) or (missing_points and len(missing_points) > 0):
            if topic not in seen:
                seen.add(topic)
                wrong_topics.append(topic)
            if len(wrong_topics) >= limit:
                break
    return wrong_topics


def _ensure_written_topic_from_input(raw_topic: str, default_subject: str, previous_wrong_topics: list[str]) -> str:
    text = (raw_topic or "").strip()
    if text:
        return text
    if previous_wrong_topics:
        return previous_wrong_topics[0]
    return f"{default_subject} topic"


def _prepare_written_question(
    exam_topic: str,
    subject: str,
    course_context: str,
    difficulty: str,
    question_type: str,
    mode: str,
    question_seed: int,
) -> dict:
    """Generate a written question from school bank first, then fallback to general generator."""
    match_result = find_relevant_school_questions(exam_topic, subject, course_context, limit=5)
    bank_item = match_result.get("best_match")
    generated = generate_written_question_from_bank(
        {
            "mode": mode,
            "question_type": question_type,
            "topic": exam_topic,
            "subject": subject,
            "difficulty": difficulty,
            "course_context": course_context,
            "best_match": bank_item,
            "question_seed": str(question_seed),
        }
    )
    if generated:
        generated["match_result"] = match_result
        generated["match_score"] = match_result.get("match_score", 0)
        generated["match_reason"] = match_result.get("match_reason", "")
        return generated

    fallback_text = (
        "Topic: {topic}\n"
        "Subject: {subject}\n"
        "Question type: {question_type}\n"
        "Course Context: {course_context}\n"
        "Generate one focused exam question in English.\n"
    ).format(
        topic=exam_topic,
        subject=subject,
        question_type=question_type,
        course_context=course_context[:6000],
    )
    generated = generate_oral_question(fallback_text, subject, difficulty)
    generated["question_source"] = "fallback_oral_llm"
    return generated


def _answer_tokens(text: str) -> set[str]:
    return {
        token
        for token in re.findall(r"[a-zA-Z][a-zA-Z0-9-]*|[0-9]+", str(text or "").lower())
        if len(token) > 2
    }


def _point_is_covered(point: str, student_answer: str) -> bool:
    point_tokens = _answer_tokens(point)
    if not point_tokens:
        return False
    answer_tokens = _answer_tokens(student_answer)
    overlap = point_tokens.intersection(answer_tokens)
    return len(overlap) >= max(1, min(3, len(point_tokens) // 2))


def _pass_level(score: float) -> str:
    if score >= 85:
        return "Excellent"
    if score >= 70:
        return "Pass"
    if score >= 50:
        return "Borderline"
    return "Needs improvement"


def grade_written_answer(params: dict) -> dict:
    question_type = str(params.get("question_type", "Short Answer 简答题"))
    student_answer = str(params.get("student_answer", "")).strip()
    expected_points = [str(item) for item in params.get("expected_points", []) if str(item).strip()]
    common_mistakes = [str(item) for item in params.get("common_mistakes", []) if str(item).strip()]
    correct_answer = str(params.get("correct_answer", "")).strip()
    model_answer = str(params.get("model_answer", "")).strip()

    if question_type == "MCQ 单选题":
        normalized_answer = student_answer.strip().upper()[:1]
        normalized_correct = correct_answer.strip().upper()[:1] or "A"
        is_correct = normalized_answer == normalized_correct
        score = 100 if is_correct else 0
        covered_points = expected_points[:2] if is_correct else []
        missing_points = [] if is_correct else expected_points[:3]
        feedback = "选择正确。" if is_correct else f"正确答案是 {normalized_correct}。"
    else:
        covered_points = [point for point in expected_points if _point_is_covered(point, student_answer)]
        missing_points = [point for point in expected_points if point not in covered_points]
        coverage = len(covered_points) / max(1, len(expected_points))
        detected_mistakes = [mistake for mistake in common_mistakes if _point_is_covered(mistake, student_answer)]
        score = max(0, min(100, round(coverage * 85 + (15 if student_answer else 0) - len(detected_mistakes) * 5, 1)))
        feedback = "答案覆盖了主要要点。" if score >= 70 else "答案还需要补充关键考点。"

    weakness_types = []
    lower_missing = " ".join(missing_points).lower()
    for marker, weakness in [
        ("definition", "definition"),
        ("classif", "classification"),
        ("mechanism", "mechanism"),
        ("pathogenesis", "mechanism"),
        ("clinical", "clinical_reasoning"),
        ("treatment", "treatment_plan"),
        ("material", "material_properties"),
        ("anatomy", "anatomy"),
        ("infection", "infection_control"),
        ("term", "terminology"),
    ]:
        if marker in lower_missing and weakness not in weakness_types:
            weakness_types.append(weakness)
    if not weakness_types and missing_points:
        weakness_types.append("exam_structure")

    return {
        "score": safe_float(score, 0),
        "pass_level": _pass_level(safe_float(score, 0) or 0),
        "level": _pass_level(safe_float(score, 0) or 0),
        "covered_points": covered_points,
        "missing_points": missing_points,
        "feedback": feedback,
        "corrected_answer": model_answer or correct_answer,
        "chinese_review": (
            f"{feedback} 已覆盖 {len(covered_points)} / {max(1, len(expected_points))} 个预期要点。"
            + (" 建议补充：" + "；".join(missing_points[:3]) if missing_points else "")
        ),
        "chinese_feedback": (
            f"{feedback} 已覆盖 {len(covered_points)} / {max(1, len(expected_points))} 个预期要点。"
            + (" 建议补充：" + "；".join(missing_points[:3]) if missing_points else "")
        ),
        "weakness_types": weakness_types,
        "next_recommended_topic": missing_points[0] if missing_points else params.get("matched_topic") or params.get("topic", ""),
    }
def render_oral_exam_mode(default_text: str):
    st.session_state.setdefault("oral_exam_history", [])
    st.session_state.setdefault("oral_exam_rounds", [])
    st.session_state.setdefault("oral_exam_rounds_target", 1)
    st.session_state.setdefault("oral_mode", WRITTEN_MODES[0])
    st.session_state.setdefault("oral_question_type", WRITTEN_QUESTION_TYPES[0])
    if st.session_state["oral_mode"] not in WRITTEN_MODES:
        st.session_state["oral_mode"] = WRITTEN_MODES[0]
    if st.session_state["oral_question_type"] not in WRITTEN_QUESTION_TYPES:
        st.session_state["oral_question_type"] = WRITTEN_QUESTION_TYPES[0]
    st.session_state.setdefault("written_question_number", 1)
    st.session_state.setdefault("oral_rounds_completed", 0)
    if "oral_student_answer" not in st.session_state:
        st.session_state["oral_student_answer"] = ""

    bank_items = load_school_question_bank()
    bank_available = bool(bank_items)
    wrong_topics = get_recent_wrong_written_topics()

    st.markdown(
        f"""
        <section class="hero">
            <div class="eyebrow">{t("written_exam_mode")}</div>
            <h1 class="hero-title">{'Generate school-style written exam questions' if get_ui_lang() == 'en' else '生成真实口试风格的笔试题'}</h1>
            <p class="hero-subtitle">{'The system prioritizes your school question bank, then uses course content to generate high-yield practice.' if get_ui_lang() == 'en' else 'AI 会优先使用学校题库命中题目，再结合课程内容生成高频考点。'}</p>
            <p class="hero-copy">{'Different modes control pressure, hints, scoring, and wrong-question review.' if get_ui_lang() == 'en' else '模式不同，难度和题目策略不同，适合日常复习、考前模拟和错题强化。'}</p>
        </section>
        """,
        unsafe_allow_html=True,
    )

    st.markdown('<div class="input-panel">', unsafe_allow_html=True)
    st.markdown(
        f'<div class="section-label">{"Written Exam Settings" if get_ui_lang() == "en" else "笔试设置"}</div>',
        unsafe_allow_html=True,
    )
    st.markdown(
        f'<div class="section-copy">{"Course content can come from Study Pack, textbooks, or lecture notes. Configure the mode first, then generate questions." if get_ui_lang() == "en" else "课程内容可来自 Study Pack、课本或课堂提要；先设置模式，再开始生成题目。"}</div>',
        unsafe_allow_html=True,
    )

    if not bank_available:
        st.warning("School question bank not found. Please add data/school_question_bank.json." if get_ui_lang() == "en" else "未找到学校题库，请先添加 data/school_question_bank.json。")

    bank_status = get_question_bank_status()
    with st.expander("调试：学校题库状态", expanded=False):
        st.markdown(f"**题库路径：** `{bank_status.get('path', 'data/school_question_bank.json')}`")
        st.markdown(f"**是否存在：** {'是' if bank_status.get('exists') else '否'}")
        st.markdown(f"**已启用题目数量：** {bank_status.get('count', 0)}")
        st.markdown(f"**状态说明：** {bank_status.get('message', '')}")

    selected_mode = st.selectbox(
        t("written_exam_mode_label"),
        WRITTEN_MODES,
        index=WRITTEN_MODES.index(st.session_state["oral_mode"]),
        format_func={
            "日常练习": "Daily Practice" if get_ui_lang() == "en" else "日常练习",
            "考前模拟": "Mock Exam" if get_ui_lang() == "en" else "考前模拟",
            "错题强化": "Wrong Question Review" if get_ui_lang() == "en" else "错题强化",
        }.get,
    )
    selected_type = st.selectbox(
        t("question_type"),
        WRITTEN_QUESTION_TYPES,
        index=WRITTEN_QUESTION_TYPES.index(st.session_state["oral_question_type"]),
        format_func={
            "自动混合": "Auto Mix" if get_ui_lang() == "en" else "自动混合",
            "MCQ 单选题": "MCQ" if get_ui_lang() == "en" else "MCQ 单选题",
            "Short Answer 简答题": "Short Answer" if get_ui_lang() == "en" else "Short Answer 简答题",
            "Case-based 病例题": "Case-based" if get_ui_lang() == "en" else "Case-based 病例题",
            "True / False 判断题": "True / False" if get_ui_lang() == "en" else "True / False 判断题",
            "Matching 匹配题": "Matching" if get_ui_lang() == "en" else "Matching 匹配题",
        }.get,
    )

    oral_default_text = st.session_state.get("last_course_text", default_text)
    oral_course_text = st.text_area(
        t("course_content"),
        value=oral_default_text,
        height=200,
        placeholder="Paste course content in English or lecture notes..." if get_ui_lang() == "en" else "粘贴课程内容（英文）或课堂提要...",
        key="oral_course_text",
    )

    st.markdown("### Configuration" if get_ui_lang() == "en" else "### 配置")
    control_col_1, control_col_2, control_col_3 = st.columns(3)
    with control_col_1:
        oral_subject = st.selectbox(
            "Subject" if get_ui_lang() == "en" else "科目",
            SUBJECT_VALUES,
            index=0,
            format_func=format_subject_option,
        )
    with control_col_2:
        oral_difficulty = st.selectbox(t("difficulty"), ["easy", "medium", "hard"], index=1)
    with control_col_3:
        default_count = 5 if selected_mode == "考前模拟" else 1
        oral_question_count = st.number_input(
            t("question_count"),
            min_value=1,
            max_value=10,
            value=default_count,
            step=1,
        )

    exam_topic = st.text_input(
        t("topic"),
        value=st.session_state.get("oral_exam_topic", ""),
        placeholder="例如：Dental caries",
        key="oral_exam_topic",
    )
    if selected_mode == "错题强化":
        inferred_topic = _ensure_written_topic_from_input(exam_topic, oral_subject, wrong_topics)
        if inferred_topic and inferred_topic != exam_topic:
            st.caption(
                f"Wrong Question Review: selected your recent weak topic: {inferred_topic}"
                if get_ui_lang() == "en"
                else f"错题强化：已为你自动选择最近错题主题「{inferred_topic}」"
            )
            exam_topic = inferred_topic
            st.session_state["oral_exam_topic"] = exam_topic

    st.session_state["oral_mode"] = selected_mode
    st.session_state["oral_question_type"] = selected_type
    st.session_state["oral_exam_rounds_target"] = int(oral_question_count) if selected_mode == "考前模拟" else 1

    if selected_mode == "错题强化" and not wrong_topics:
        st.info("No wrong-question records yet. Please complete Daily Practice or Mock Exam first." if get_ui_lang() == "en" else "你还没有错题记录，请先完成日常练习或考前模拟。")

    if st.button(t("generate_written_question"), type="primary", use_container_width=True):
        if not oral_course_text.strip():
            st.error("Please enter course content first." if get_ui_lang() == "en" else "请先填写课程内容。")
        elif selected_mode == "错题强化" and not wrong_topics:
            st.info("No wrong-question records yet. Please complete Daily Practice or Mock Exam first." if get_ui_lang() == "en" else "你还没有错题记录，请先完成日常练习或考前模拟。")
        else:
            try:
                prepared = run_with_dental_loader(
                    "Generating question..." if get_ui_lang() == "en" else "正在生成题目...",
                    lambda: _prepare_written_question(
                        exam_topic=exam_topic.strip() or _ensure_written_topic_from_input("", oral_subject, wrong_topics),
                        subject=oral_subject,
                        course_context=oral_course_text,
                        difficulty=oral_difficulty,
                        question_type=selected_type,
                        mode=selected_mode,
                        question_seed=int(time.time()),
                    ),
                    "AI is preparing your written practice question." if get_ui_lang() == "en" else "AI 正在准备题目，请稍等。",
                )
                prepared["question_type"] = selected_type
                prepared["topic"] = prepared.get("topic") or exam_topic.strip() or oral_subject
                prepared["difficulty"] = prepared.get("difficulty") or oral_difficulty
                st.session_state["oral_question_data"] = prepared
                st.session_state["oral_exam_result"] = None
                st.session_state["reset_oral_student_answer"] = True
                if selected_mode == "考前模拟":
                    st.session_state["oral_exam_rounds"] = []
                else:
                    st.session_state["oral_exam_rounds"] = st.session_state.get("oral_exam_rounds", [])[:0]
                st.session_state["oral_rounds_completed"] = 0
                st.session_state["last_course_text"] = oral_course_text
                st.success("已生成题目。")
            except Exception as exc:
                st.error(f"生成题目失败：{exc}")

    question_data = st.session_state.get("oral_question_data")
    if question_data:
        question_source = question_data.get("question_source") or question_data.get("source") or question_data.get("bank_source", "")
        source_text = "学校题库" if question_source in ("school_question_bank", "school_question_bank") else "语言模型生成"
        st.markdown("### 题目")
        st.info(question_data.get("question", ""))
        if question_data.get("options"):
            for option in question_data.get("options", []):
                st.markdown(f"- {option}")
        with st.expander("题目信息与评分要点"):
            st.markdown(f"**题库来源：** {source_text}")
            st.markdown(f"**题型：** {question_data.get('question_type', selected_type)}")
            st.markdown(f"**难度：** {question_data.get('difficulty', oral_difficulty)}")
            st.markdown(f"**主题：** {question_data.get('matched_topic') or question_data.get('topic', exam_topic)}")
            if question_data.get("focus") or question_data.get("question_focus"):
                st.markdown(f"**Focus：** {question_data.get('focus') or question_data.get('question_focus')}")
            st.markdown(f"**匹配分数：** {question_data.get('match_score', 0)}")
            st.markdown(f"**匹配原因：** {question_data.get('match_reason', '')}")
            can_show_hints = selected_mode != "考前模拟" or bool(st.session_state.get("oral_exam_result"))
            if selected_mode == "考前模拟" and not can_show_hints:
                st.caption("考前模拟模式下，作答前不显示提示和评分细则。提交后可以查看解析。")
            if can_show_hints and question_data.get("expected_points"):
                st.markdown("**Must-know 要点：**")
                for item in question_data.get("expected_points", []):
                    st.markdown(f"- {item}")
            if can_show_hints and question_data.get("common_mistakes"):
                st.markdown("**常见扣分点：**")
                for item in question_data.get("common_mistakes", []):
                    st.markdown(f"- {item}")
            if can_show_hints and question_data.get("must_mention_terms"):
                st.markdown("**提示词：**")
                for item in question_data.get("must_mention_terms", [])[:5]:
                    st.markdown(f"- {item}")
            if can_show_hints and question_data.get("follow_up_questions"):
                st.markdown("**可考察追问方向：**")
                for item in question_data.get("follow_up_questions", [])[:4]:
                    st.markdown(f"- {item}")
            if can_show_hints and question_data.get("scoring_rubric"):
                st.markdown("**评分 Rubric：**")
                st.json(question_data.get("scoring_rubric"))

        with st.expander("调试：笔试题库匹配"):
            st.write(
                {
                    "matched_topic": question_data.get("matched_topic") or question_data.get("topic"),
                    "match_score": question_data.get("match_score", 0),
                    "inferred_focus": question_data.get("focus") or question_data.get("question_focus"),
                    "question_source": question_source,
                    "generated_question": question_data.get("question"),
                }
            )

        if st.session_state.get("reset_oral_student_answer"):
            st.session_state["oral_student_answer"] = ""
            st.session_state["reset_oral_student_answer"] = False

        student_answer = st.text_area(
            "我的答案",
            height=180,
            placeholder="Type your written exam answer in English...",
            key="oral_student_answer",
        )

        if st.button(t("submit_answer"), type="primary", use_container_width=True):
            if not student_answer.strip():
                st.error("请先输入你的答案。")
            else:
                try:
                    result = run_with_dental_loader(
                        "正在评卷..." if get_ui_lang() != "en" else "Grading answer...",
                        lambda: grade_written_answer(
                            {
                                "question_type": question_data.get("question_type", selected_type),
                                "question": question_data.get("question", ""),
                                "options": question_data.get("options", []),
                                "correct_answer": question_data.get("correct_answer", ""),
                                "student_answer": student_answer,
                                "expected_points": question_data.get("expected_points", []),
                                "model_answer": question_data.get("model_answer", ""),
                                "scoring_rubric": question_data.get("scoring_rubric", {}),
                                "common_mistakes": question_data.get("common_mistakes", []),
                                "mode": selected_mode,
                                "topic": question_data.get("matched_topic") or question_data.get("topic", exam_topic),
                            }
                        ),
                        "AI 教授正在分析你的答案，请稍候..." if get_ui_lang() != "en" else "AI is analyzing your answer...",
                    )
                    st.session_state["oral_exam_result"] = result
                    attempt = {
                        "session_id": st.session_state.get("oral_session_id"),
                        "subject": oral_subject,
                        "difficulty": question_data.get("difficulty", oral_difficulty),
                        "topic": question_data.get("topic", exam_topic),
                        "course_context": oral_course_text,
                        "question": question_data.get("question", ""),
                        "answer": student_answer,
                        "result": result,
                        "grading_result": result,
                        "expected_points": question_data.get("expected_points", []),
                        "model_answer": question_data.get("model_answer", ""),
                        "options": question_data.get("options", []),
                        "correct_answer": question_data.get("correct_answer", ""),
                        "question_type": selected_type,
                        "question_source": question_source,
                    }
                    try:
                        save_written_exam_attempt(attempt)
                        update_user_weaknesses_from_attempt(
                            oral_subject,
                            attempt.get("topic") or oral_subject,
                            result.get("missing_points") or [],
                            result.get("score"),
                        )
                        increment_usage("text_exam", 1)
                    except Exception as exc:
                        st.warning("笔试记录保存失败，请稍后重试。")
                        if st.secrets.get("STAGE", "").lower() == "dev" or os.getenv("DENTPILOT_DEBUG", "").lower() in {"1", "true", "yes", "on"}:
                            with st.expander("调试信息"):
                                st.caption(f"{type(exc).__name__}: {exc}")

                    st.session_state["oral_exam_history"].insert(0, attempt)
                    st.session_state["oral_exam_history"] = st.session_state["oral_exam_history"][:10]
                    st.session_state["oral_exam_rounds"].append(attempt)
                    st.session_state["oral_rounds_completed"] = len(st.session_state["oral_exam_rounds"])
                except OralExamConfigError as exc:
                    st.error(str(exc))
                except OralExamJSONError as exc:
                    st.error("DeepSeek 返回异常，请检查输出格式。")
                    st.code(exc.raw_output)
                except Exception as exc:
                    st.error(f"评分失败：{exc}")

    if st.session_state.get("oral_exam_result"):
        st.markdown("### 口试解析")
        render_oral_grade_result(st.session_state["oral_exam_result"])

        completed = len(st.session_state.get("oral_exam_rounds", []))
        target = st.session_state.get("oral_exam_rounds_target", 1)
        if completed < target:
            if st.button(t("practice_similar"), use_container_width=True):
                st.session_state["oral_exam_result"] = None
                st.session_state["oral_question_data"] = None
                st.session_state["reset_oral_student_answer"] = True
                st.rerun()
        elif selected_mode == "考前模拟":
            st.success(f"本轮已完成 {completed} 题，可切换主题或题型继续下一轮。")

    if selected_mode == "考前模拟":
        st.markdown("### 考前模拟进度")
        st.progress(min(len(st.session_state.get("oral_exam_rounds", [])), st.session_state.get("oral_exam_rounds_target", 1)) / max(1, st.session_state.get("oral_exam_rounds_target", 1)))
        st.caption(f"已做：{len(st.session_state.get('oral_exam_rounds', []))} / {st.session_state.get('oral_exam_rounds_target', 1)}")

    st.markdown("</div>", unsafe_allow_html=True)

    rounds = st.session_state.get("oral_exam_rounds", [])
    if rounds:
        st.markdown("### 当前口试 Session")
        st.caption(f"本次已完成 {len(rounds)} 题。")
        for index, attempt in enumerate(rounds, start=1):
            result = attempt.get("grading_result") or attempt.get("result", {})
            with st.expander(f"题 {index}: {result.get('score', 0)}/100 ({result.get('level', '')})"):
                st.markdown(f"**题目：** {attempt.get('question', '')}")
                st.markdown(f"**我的答案：** {attempt.get('answer', '')}")
                st.markdown(f"**反馈：** {result.get('chinese_feedback', '')}")

    st.markdown(f"### {t('recent_written_records')}")
    if st.session_state.get("oral_exam_history_error"):
        st.error(f"读取笔试记录失败：{st.session_state['oral_exam_history_error']}")
    history = st.session_state.get("oral_exam_history", [])
    if not history:
        st.info("还没有笔试记录。请先生成题目并提交。")
    else:
        for index, attempt in enumerate(history[:5], start=1):
            result = attempt.get("result", {})
            score = result.get("score", attempt.get("score", 0))
            level = result.get("level", "")
            label = f"{index}. {attempt.get('topic') or attempt.get('subject')} - {score}/100 ({level})"
            with st.expander(label):
                st.markdown(f"**题目：** {attempt.get('question', '')}")
                st.markdown(f"**我的答案：** {attempt.get('answer') or attempt.get('student_answer', '')}")
                st.markdown(f"**反馈：** {result.get('chinese_feedback') or attempt.get('feedback', '')}")

    wrong_records = []
    for attempt in history:
        result = attempt.get("result") or attempt.get("grading_result") or attempt
        score = safe_float(result.get("score"), None)
        missing_points = result.get("missing_points") or attempt.get("missing_points") or []
        is_wrong = (score is not None and score < 70) or bool(missing_points)
        if is_wrong:
            wrong_records.append(attempt)

    st.markdown(f"### {t('wrong_question_notebook')}")
    if not wrong_records:
        st.info("暂时没有错题记录。完成日常练习或考前模拟后，低分题和缺失要点会出现在这里。")
    else:
        for index, attempt in enumerate(wrong_records[:8], start=1):
            result = attempt.get("result") or attempt.get("grading_result") or attempt
            score = safe_float(result.get("score"), 0)
            topic_label = attempt.get("topic") or attempt.get("subject") or "Written Exam"
            with st.expander(f"{index}. {topic_label} - {score}/100"):
                st.markdown(f"**题目：** {attempt.get('question', '')}")
                st.markdown(f"**我的答案：** {attempt.get('answer') or attempt.get('student_answer', '')}")
                model_answer = attempt.get("model_answer") or result.get("corrected_answer") or ""
                if model_answer:
                    st.markdown(f"**参考答案：** {model_answer}")
                missing_points = result.get("missing_points") or attempt.get("missing_points") or []
                if missing_points:
                    st.markdown("**缺失要点：**")
                    for point in missing_points[:5]:
                        st.markdown(f"- {point}")
                if st.button(t("practice_similar"), key=f"retry_written_{index}_{attempt.get('id', '')}"):
                    st.session_state["oral_exam_topic"] = topic_label
                    st.session_state["oral_mode"] = "错题强化"
                    st.session_state["oral_question_data"] = None
                    st.session_state["oral_exam_result"] = None
                    st.session_state["reset_oral_student_answer"] = True
                    st.rerun()

    st.markdown(
        """
        <div class="footer">
            本模块支持中文界面，按“笔试模式/题型”控制题目来源和难度。Study Pack、临床病例与实时口试功能可持续使用。
        </div>
        """,
        unsafe_allow_html=True,
    )
def render_clinical_case_grade(result: dict):
    st.markdown("### 病例评分")
    score_col, level_col = st.columns(2)
    score_col.metric("得分", f"{result.get('score', 0)}/100")
    level_col.metric("等级", result.get("level", ""))

    rubric_rows = [
        {"维度": "诊断", "得分": result.get("diagnosis_score", 0), "权重": 20},
        {"维度": "依据与证据", "得分": result.get("evidence_score", 0), "权重": 20},
        {"维度": "鉴别诊断", "得分": result.get("differential_score", 0), "权重": 15},
        {"维度": "额外检查", "得分": result.get("tests_score", 0), "权重": 15},
        {"维度": "治疗计划", "得分": result.get("treatment_score", 0), "权重": 15},
        {"维度": "医患沟通", "得分": result.get("patient_communication_score", 0), "权重": 10},
        {"维度": "安全与红旗", "得分": result.get("safety_score", 0), "权重": 5},
    ]
    st.dataframe(pd.DataFrame(rubric_rows), use_container_width=True, hide_index=True)

    col_1, col_2 = st.columns(2)
    with col_1:
        st.subheader("做得好的点")
        strengths = result.get("strengths") or []
        if strengths:
            for item in strengths:
                st.markdown(f"- {item}")
        else:
            st.write("未识别到明显优点")
    with col_2:
        st.subheader("待改进")
        missing_points = result.get("missing_points") or []
        if missing_points:
            for item in missing_points:
                st.markdown(f"- {item}")
        else:
            st.write("未识别到明显缺失点")

    st.subheader("Model Answer")
    st.write(result.get("model_answer", ""))

    st.subheader("中文反馈")
    st.write(result.get("chinese_feedback", ""))

    st.subheader("下一步建议")
    st.info(result.get("next_practice_suggestion", ""))


def render_clinical_case_mode(default_text: str):
    st.session_state.setdefault("clinical_case_history", [])

    st.markdown(
        f"""
        <section class="hero">
            <div class="eyebrow">{t("clinical_case_training")}</div>
            <h1 class="hero-title">{t("clinical_case_training")}</h1>
            <p class="hero-subtitle">根据课程内容生成牙科/医学临床病例，训练诊断、证据、鉴别诊断、检查、治疗计划和患者沟通。</p>
            <p class="hero-copy">完成病例训练后，系统会保存记录并用于后续弱点分析。</p>
        </section>
        """,
        unsafe_allow_html=True,
    )

    st.markdown('<div class="input-panel">', unsafe_allow_html=True)
    st.markdown(f'<div class="section-label">{t("case_materials")}</div>', unsafe_allow_html=True)
    st.markdown(
        '<div class="section-copy">粘贴英文课程内容，系统会生成一个虚构但贴近课程重点的临床病例。</div>',
        unsafe_allow_html=True,
    )

    case_default_text = st.session_state.get("last_course_text", default_text)
    case_course_text = st.text_area(
        t("course_content"),
        value=case_default_text,
        height=220,
        placeholder="Paste your English dental or medical course text here...",
        key="clinical_case_course_text",
    )

    control_col_1, control_col_2 = st.columns(2)
    with control_col_1:
        case_subject = st.selectbox(
            "Subject" if get_ui_lang() == "en" else "科目",
            SUBJECT_VALUES,
            key="clinical_case_subject",
            format_func=format_subject_option,
        )
    with control_col_2:
        case_difficulty = st.selectbox(
            "Difficulty",
            ["easy", "medium", "hard"],
            index=1,
            key="clinical_case_difficulty",
        )

    if st.button(t("generate_case"), type="primary", use_container_width=True):
        if not case_course_text.strip():
            st.error("请先粘贴课件内容。")
        else:
            try:
                st.session_state["clinical_case_data"] = run_with_dental_loader(
                    "正在生成病例题..." if get_ui_lang() != "en" else "Generating case questions...",
                    lambda: generate_clinical_case(
                        case_course_text,
                        case_subject,
                        case_difficulty,
                    ),
                    "Clinical case materials are being generated from your course text." if get_ui_lang() == "en" else "根据课程内容生成病例材料，请稍候。",
                )
                st.session_state["clinical_case_result"] = None
                st.session_state["last_course_text"] = case_course_text
                st.success("病例已生成。")
            except ClinicalCaseConfigError as exc:
                st.error(str(exc))
            except ClinicalCaseJSONError as exc:
                st.error("DeepSeek 返回了不可解析的 JSON，请检查输出格式。")
                st.code(exc.raw_output)
            except Exception as exc:
                st.error(f"生成病例失败：{exc}")

    case_data = st.session_state.get("clinical_case_data")
    if case_data:
        st.markdown(f"### {t('case_materials')}")
        st.subheader(case_data.get("case_title", "Clinical Case"))

        case_col_1, case_col_2 = st.columns(2)
        with case_col_1:
            st.markdown("**Patient Info**")
            st.write(case_data.get("patient_info", ""))
            st.markdown("**Chief Complaint**")
            st.write(case_data.get("chief_complaint", ""))
            st.markdown("**History**")
            st.write(case_data.get("history", ""))
        with case_col_2:
            st.markdown("**Clinical Findings**")
            st.write(case_data.get("clinical_findings", ""))
            st.markdown("**Radiographic Findings**")
            st.write(case_data.get("radiographic_findings", ""))

        st.markdown("### Questions")
        for index, question in enumerate(case_data.get("questions", []), start=1):
            st.markdown(f"{index}. {question}")

        with st.expander("鏌ョ湅 expected diagnosis / expected points / red flags"):
            st.markdown("**Expected Diagnosis**")
            st.write(case_data.get("expected_diagnosis", ""))
            st.markdown("**Expected Points**")
            for item in case_data.get("expected_points", []):
                st.markdown(f"- {item}")
            st.markdown("**Red Flags**")
            for item in case_data.get("red_flags", []):
                st.markdown(f"- {item}")

        student_answer = st.text_area(
            t("my_analysis"),
            height=220,
            placeholder=(
                "Answer in English. You can organize it as: diagnosis, evidence, "
                "differentials, tests, treatment plan, patient explanation."
            ),
            key="clinical_case_answer",
        )

        if st.button("Submit Case Answer", type="primary", use_container_width=True):
            if not student_answer.strip():
                st.error("请先输入你的病例分析答案。")
            else:
                try:
                    result = run_with_dental_loader(
                        "正在批改病例分析..." if get_ui_lang() != "en" else "Reviewing case analysis...",
                        lambda: grade_clinical_case(case_data, student_answer, case_subject),
                        "病例分析评分需要一些时间..." if get_ui_lang() != "en" else "Case grading in progress...",
                    )
                    st.session_state["clinical_case_result"] = result
                    attempt = {
                        "subject": case_subject,
                        "difficulty": case_difficulty,
                        "case_title": case_data.get("case_title", ""),
                        "case_data": case_data,
                        "answer": student_answer,
                        "result": result,
                    }
                    try:
                        save_clinical_case_attempt(attempt)
                        update_user_weaknesses_from_attempt(
                            case_subject,
                            case_data.get("case_title", case_subject),
                            result.get("missing_points") or [],
                            result.get("score"),
                        )
                        increment_usage("case", 1)
                    except Exception as exc:
                        st.warning(f"病例记录保存失败：{exc}")
                    st.session_state["clinical_case_history"].insert(0, attempt)
                    st.session_state["clinical_case_history"] = st.session_state["clinical_case_history"][:10]
                except ClinicalCaseConfigError as exc:
                    st.error(str(exc))
                except ClinicalCaseJSONError as exc:
                    st.error("DeepSeek 返回了无效 JSON。下面是原始输出，方便调试：")
                    st.code(exc.raw_output)
                except Exception as exc:
                    st.error(f"评改失败：{exc}")

    if st.session_state.get("clinical_case_result"):
        render_clinical_case_grade(st.session_state["clinical_case_result"])

    st.markdown("</div>", unsafe_allow_html=True)

    st.markdown("### 病例历史")
    if st.session_state.get("clinical_case_history_error"):
        st.error(f"读取病例记录失败：{st.session_state['clinical_case_history_error']}")
    history = st.session_state.get("clinical_case_history", [])
    if not history:
        st.info("暂无历史病例记录，完成一次分析后会记录在这里。")
    else:
        for index, attempt in enumerate(history[:5], start=1):
            result = attempt.get("result", {})
            score = result.get("score", attempt.get("score", 0))
            level = result.get("level", "")
            label = f"{index}. {attempt.get('case_title') or attempt.get('subject')} - {score}/100 ({level})"
            with st.expander(label):
                st.markdown(f"**Subject:** {attempt.get('subject', '')}")
                st.markdown(f"**Your Answer:** {attempt.get('answer') or attempt.get('student_answer', '')}")
                st.markdown(f"**Chinese Feedback:** {result.get('chinese_feedback') or attempt.get('feedback', '')}")

    st.markdown(
        """
        <div class="footer">
            Clinical Case Training 是 DentPilot AI 的文本病例训练模式，用于练习临床思维，不替代真实医疗建议。
        </div>
        """,
        unsafe_allow_html=True,
    )


def render_weakness_analysis_mode():
    oral_history = st.session_state.get("oral_exam_history", [])
    clinical_history = st.session_state.get("clinical_case_history", [])
    realtime_oral_history = st.session_state.get("realtime_oral_history", [])
    study_pack_records = st.session_state.get("study_pack_records", [])
    combined_exam_history = [*oral_history, *realtime_oral_history]

    st.markdown(
        f"""
        <section class="hero">
            <div class="eyebrow">{t("weakness_analysis")}</div>
            <h1 class="hero-title">{t("weakness_analysis")}</h1>
            <p class="hero-subtitle">{"Analyze written exam, oral exam, and clinical case records to identify strengths, weak points, and next practice suggestions." if get_ui_lang() == "en" else "根据笔试、口试和临床病例训练记录，找出强项、弱点和可能原因，并生成 3 天复习计划。"}</p>
            <p class="hero-copy">{"The system will summarize your weaknesses after practice." if get_ui_lang() == "en" else "完成训练后，系统会分析你的薄弱知识点，并生成个性化复习计划。"}</p>
        </section>
        """,
        unsafe_allow_html=True,
    )

    st.markdown('<div class="input-panel">', unsafe_allow_html=True)
    st.markdown('<div class="section-label">练习记录</div>', unsafe_allow_html=True)

    metric_col_1, metric_col_2, metric_col_3, metric_col_4 = st.columns(4)
    metric_col_1.metric("Written Exam Attempts" if get_ui_lang() == "en" else "笔试记录", len(oral_history))
    metric_col_2.metric("Clinical Case Attempts" if get_ui_lang() == "en" else "病例记录", len(clinical_history))
    metric_col_3.metric("Oral Exam Attempts" if get_ui_lang() == "en" else "口试记录", len(realtime_oral_history))
    metric_col_4.metric("Study Packs" if get_ui_lang() == "en" else "复习包", len(study_pack_records))

    if not combined_exam_history and not clinical_history:
        st.info("请先完成至少 1 次口试或笔试后再分析。")
        st.markdown("</div>", unsafe_allow_html=True)
        return

    if st.button("Analyze My Weaknesses", type="primary", use_container_width=True):
        try:
            st.session_state["weakness_analysis_result"] = run_with_dental_loader(
                "正在分析你的练习弱点..." if get_ui_lang() != "en" else "Analyzing your practice weaknesses...",
                lambda: analyze_weaknesses(
                    combined_exam_history,
                    clinical_history,
                ),
                "系统正在整理错题画像与复习建议..." if get_ui_lang() != "en" else "Generating weakness summary and recommendations.",
            )
            for attempt in [*combined_exam_history, *clinical_history]:
                result_data = attempt.get("result") or attempt
                update_user_weaknesses_from_attempt(
                    attempt.get("subject") or "Dentistry",
                    attempt.get("topic") or attempt.get("case_title") or "General practice",
                    result_data.get("missing_points") or [],
                    result_data.get("score"),
                )
            st.success("弱项分析已完成。")
        except WeaknessAnalysisConfigError as exc:
            st.error(str(exc))
        except WeaknessAnalysisJSONError as exc:
            st.error("DeepSeek 返回了无效 JSON。下面是原始输出，方便调试：")
            st.code(exc.raw_output)
        except Exception as exc:
            st.error(f"弱项分析失败：{exc}")

    result = st.session_state.get("weakness_analysis_result")
    if result:
        st.markdown("### Overall Summary")
        st.write(result.get("overall_summary", ""))

        col_1, col_2 = st.columns(2)
        with col_1:
            st.subheader("Strong Areas")
            for item in result.get("strong_areas", []):
                st.markdown(f"- {item}")
        with col_2:
            st.subheader("Weak Areas")
            for item in result.get("weak_areas", []):
                st.markdown(f"- {item}")

        st.subheader("Likely Reasons")
        for item in result.get("likely_reasons", []):
            st.markdown(f"- {item}")

        st.subheader("Topic Breakdown")
        topic_breakdown = result.get("topic_breakdown", [])
        if topic_breakdown:
            st.dataframe(pd.DataFrame(topic_breakdown), use_container_width=True, hide_index=True)
        else:
            st.info("暂无 Topic Breakdown。")

        st.subheader("Three-Day Study Plan")
        for day_plan in result.get("three_day_plan", []):
            with st.expander(f"Day {day_plan.get('day')}: {day_plan.get('focus', '')}", expanded=True):
                for task in day_plan.get("tasks", []):
                    st.markdown(f"- {task}")

        with st.expander("Next Practice Suggestions"):
            st.markdown("**Next Oral Questions**")
            for item in result.get("next_oral_questions", []):
                st.markdown(f"- {item}")
            st.markdown("**Next Case Topics**")
            for item in result.get("next_case_topics", []):
                st.markdown(f"- {item}")

    st.markdown("</div>", unsafe_allow_html=True)

    st.markdown(
        """
        <div class="footer">
            Weakness Analysis 会根据当前用户的练习记录生成学习建议；刷新后仍会优先读取 Supabase 中的历史记录。
        </div>
        """,
        unsafe_allow_html=True,
    )


with st.sidebar:
    render_ui_language_selector()
    st.markdown("---")
    render_sidebar_account()
    render_admin_dashboard()
    st.markdown("## DentPilot AI")
    st.caption(
        "AI Study Assistant for English-Taught Dental/Medical Students"
        if get_ui_lang() == "en"
        else "面向中国留学生的英授牙科/医学学习助手"
    )
    st.markdown("---")

    mode_options = ["Study Pack", "AI Written Exam", "Clinical Case", "Weakness Analysis", "Realtime Oral Exam"]
    selected_mode_key = load_selected_mode()
    selected_mode_label = ALLOWED_MODE_KEYS.get(selected_mode_key, "Study Pack")
    mode = st.radio(
        t("study_mode"),
        mode_options,
        index=mode_options.index(selected_mode_label) if selected_mode_label in mode_options else 0,
        format_func={
            "Study Pack": t("study_pack_mode"),
            "AI Written Exam": t("written_exam_mode"),
            "Clinical Case": t("clinical_case_mode"),
            "Weakness Analysis": t("weakness_analysis_mode"),
            "Realtime Oral Exam": t("realtime_oral_exam_mode"),
        }.get,
    )
    save_selected_mode(MODE_LABEL_TO_KEY.get(mode, "study_pack"))

    mode_descriptions = {
        "Study Pack": "Upload course materials and generate study packs, quizzes, and Anki cards." if get_ui_lang() == "en" else "上传课程材料，自动生成复习包、Quiz 和 Anki 卡片。",
        "AI Written Exam": "Generate school-style written exam questions with scoring and feedback." if get_ui_lang() == "en" else "按模式/题型生成学校风格写作题，并给出详细评分与错题反馈。",
        "Clinical Case": "Practice clinical diagnosis and treatment planning with case questions." if get_ui_lang() == "en" else "基于案例的口腔临床诊断与治疗问答训练。",
        "Weakness Analysis": "Analyze practice history and generate short-term revision suggestions." if get_ui_lang() == "en" else "基于历史记录分析薄弱环节，给出短期训练建议。",
        "Realtime Oral Exam": "Open the standalone realtime oral exam room." if get_ui_lang() == "en" else "实时口语练习与对话反馈（独立应用会话入口）。",
    }

    st.markdown("### Current Mode" if get_ui_lang() == "en" else "### 当前模式")
    st.write(mode_descriptions.get(mode, "DentPilot AI 功能模块"))

    if mode != "Realtime Oral Exam":
        st.markdown("---")
        if os.getenv("DEEPSEEK_API_KEY"):
            st.caption("AI 服务配置：已启用")
        else:
            st.caption("AI 服务未配置：请在环境变量中设置 DeepSeek API Key。")


    st.markdown("---")
    render_my_learning_summary()


if mode == "AI Written Exam":
    render_oral_exam_mode(sample)
    st.stop()

if mode == "Clinical Case":
    render_clinical_case_mode(sample)
    st.stop()

if mode == "Weakness Analysis":
    render_weakness_analysis_mode()
    st.stop()

if mode == "Realtime Oral Exam":
    render_realtime_oral_exam_page(get_ui_lang())
    st.markdown("### Recent Realtime Oral Exam Records" if get_ui_lang() == "en" else "### 最近实时口试记录")
    if st.session_state.get("realtime_oral_history_error"):
        st.error(f"读取实时口试记录失败：{st.session_state['realtime_oral_history_error']}")
    try:
        usage = get_usage_today() or {}
        st.info(f"今日语音用量：{float(usage.get('voice_minutes_used') or 0):.1f} 分钟")
    except Exception as exc:
        st.warning(f"语音用量读取失败：{exc}")
    realtime_history = st.session_state.get("realtime_oral_history", [])
    if not realtime_history:
        st.info("还没有实时口试记录。实时口试记录由新版口试 App 保存；这里不会显示虚假记录。")
    else:
        for index, attempt in enumerate(realtime_history[:5], start=1):
            label = f"{index}. {attempt.get('topic') or attempt.get('subject') or 'Oral exam'} - {attempt.get('score', '暂无')}"
            with st.expander(label):
                st.markdown(f"**Question:** {attempt.get('question', '')}")
                st.markdown(f"**Your Answer:** {attempt.get('student_answer', '')}")
                st.markdown(f"**Feedback:** {attempt.get('feedback', '')}")
    st.stop()


st.markdown(
    f"""
    <section class="hero">
        <div class="eyebrow">{t("home_eyebrow")}</div>
        <h1 class="hero-title">DentPilot AI</h1>
        <p class="hero-subtitle">{t("home_subtitle")}</p>
        <p class="hero-copy">{t("home_copy")}</p>
        <div class="hero-metrics">
            <div class="metric-pill">
                <div class="metric-value">{t("metric_cn_value")}</div>
                <div class="metric-label">{t("metric_cn_label")}</div>
            </div>
            <div class="metric-pill">
                <div class="metric-value">{t("metric_quiz_value")}</div>
                <div class="metric-label">{t("metric_quiz_label")}</div>
            </div>
            <div class="metric-pill">
                <div class="metric-value">Anki</div>
                <div class="metric-label">{t("metric_anki_label")}</div>
            </div>
        </div>
    </section>
    """,
    unsafe_allow_html=True,
)


st.markdown(
    f"""
    <div class="feature-grid">
        <div class="feature-card">
            <div class="feature-icon">01</div>
            <div class="feature-title">{t("feature_cn_title")}</div>
            <div class="feature-copy">{t("feature_cn_copy")}</div>
        </div>
        <div class="feature-card">
            <div class="feature-icon">02</div>
            <div class="feature-title">{t("feature_terms_title")}</div>
            <div class="feature-copy">{t("feature_terms_copy")}</div>
        </div>
        <div class="feature-card">
            <div class="feature-icon">03</div>
            <div class="feature-title">{t("feature_quiz_title")}</div>
            <div class="feature-copy">{t("feature_quiz_copy")}</div>
        </div>
        <div class="feature-card">
            <div class="feature-icon">04</div>
            <div class="feature-title">{t("feature_export_title")}</div>
            <div class="feature-copy">{t("feature_export_copy")}</div>
        </div>
    </div>
    """,
    unsafe_allow_html=True,
)


st.markdown('<div class="input-panel">', unsafe_allow_html=True)
st.markdown(f'<div class="section-label">{t("study_pack_upload")}</div>', unsafe_allow_html=True)
st.markdown(
    f'<div class="section-copy">{t("study_pack_support")}</div>',
    unsafe_allow_html=True,
)

st.markdown(
    f"""
    <div class="example-grid">
        <div class="example-card">
            <div class="example-title">{t("example_input")}</div>
            <div class="example-copy">{sample}</div>
        </div>
        <div class="workflow-card">
            <div class="example-title">{t("generated_content")}</div>
            <div class="example-copy">{t("study_pack_workflow")}</div>
        </div>
    </div>
    """,
    unsafe_allow_html=True,
)

uploaded_course_file = st.file_uploader(
    t("upload_course_file"),
    type=["pdf", "docx", "pptx", "ppt", "txt"],
    help=t("upload_course_help"),
)

uploaded_text = ""
extraction_report = None
uploaded_source_name = ""
uploaded_page_or_slide_count = 0
if uploaded_course_file is not None:
    try:
        uploaded_source_name = uploaded_course_file.name
        file_suffix = Path(uploaded_course_file.name).suffix.lower()
        if file_suffix == ".pdf":
            uploaded_text, extraction_report = run_with_dental_loader(
                "正在分析课件内容，请稍等..." if get_ui_lang() == "en" else "正在分析课件内容，请稍等...",
                lambda: extract_pdf_text(uploaded_course_file),
                "大文件会分块处理，请耐心等待。" if get_ui_lang() == "en" else "大文件会分块处理，请稍候。",
            )
            source_label = f"PDF 的 {extraction_report.get('total_pages', 0)} 页"
            uploaded_page_or_slide_count = int(extraction_report.get("total_pages", 0))
        elif file_suffix == ".docx":
            uploaded_text, paragraph_count = run_with_dental_loader(
                "正在分析课件内容，请稍等..." if get_ui_lang() == "en" else "正在分析课件内容，请稍等...",
                lambda: extract_docx_text(uploaded_course_file),
                "正在读取 Word 文本..." if get_ui_lang() == "en" else "正在读取 Word 文本...",
            )
            source_label = f"Word 文档的 {paragraph_count} 个段落"
            uploaded_page_or_slide_count = int(paragraph_count)
        elif file_suffix == ".pptx":
            uploaded_text, ppt_report = run_with_dental_loader(
                "正在分析课件内容，请稍等..." if get_ui_lang() == "en" else "正在分析课件内容，请稍等...",
                lambda: extract_pptx_text(uploaded_course_file),
                "正在读取 PPT 内容，请稍候。" if get_ui_lang() == "en" else "正在读取 PPT 内容，请稍候。",
            )
            source_label = f"PowerPoint 的 {ppt_report.get('total_slides', 0)} 页幻灯片"
            uploaded_page_or_slide_count = int(ppt_report.get("total_slides", 0))
            with st.expander("PowerPoint 提取报告", expanded=True):
                st.markdown(f"**文件名：** {ppt_report.get('filename', '')}")
                st.markdown(f"**总幻灯片数：** {ppt_report.get('total_slides', 0)}")
                st.markdown(f"**成功提取幻灯片数：** {ppt_report.get('successful_slides', 0)}")
                st.markdown(f"**提取字符数：** {ppt_report.get('char_count', 0)}")
                st.text_area(
                    "前 1000 字预览",
                    value=str(ppt_report.get("preview", "")),
                    height=220,
                    disabled=True,
                    key="pptx_extract_preview",
                )
                total_slides = max(1, int(ppt_report.get("total_slides", 1)))
                char_count = int(ppt_report.get("char_count", 0))
                if char_count < 500 or char_count / total_slides < 80:
                    st.warning("PPT 可能主要是图片或扫描内容，当前只能读取可选中文本。请导出为可复制文字的 PPTX/PDF，或复制粘贴讲义文本。")
        elif file_suffix == ".ppt":
            source_label = "旧版 PowerPoint 文件"
            st.warning("暂不支持直接解析旧版 .ppt 文件。请先另存为 .pptx 后再上传，或复制粘贴幻灯片文字。")
        elif file_suffix == ".txt":
            uploaded_text, line_count = run_with_dental_loader(
                "正在分析课件内容，请稍等..." if get_ui_lang() == "en" else "正在分析课件内容，请稍等...",
                lambda: extract_txt_text(uploaded_course_file),
                "正在读取 TXT 内容..." if get_ui_lang() == "en" else "正在读取 TXT 内容...",
            )
            source_label = f"TXT 文件的 {line_count} 行"
        else:
            source_label = "课程文档"
            st.warning("暂时只支持 PDF、Word(docx) 和 TXT 文件。")
        if uploaded_text.strip():
            st.success(f"已从 {source_label} 中提取文本。生成前你仍然可以编辑。")
            if extraction_report:
                with st.expander("PDF 提取报告", expanded=True):
                    st.markdown(f"**文件名：** {extraction_report.get('filename', '')}")
                    st.markdown(f"**总页数：** {extraction_report.get('total_pages', 0)}")
                    st.markdown(f"**成功提取页数：** {extraction_report.get('successful_pages', 0)}")
                    st.markdown(f"**提取字符数：** {extraction_report.get('char_count', 0)}")
                    if extraction_report.get("failed_pages"):
                        st.warning(f"提取失败页：{extraction_report.get('failed_pages')}")
                    if extraction_report.get("empty_pages"):
                        st.caption(f"空白或无可选中文本页：{extraction_report.get('empty_pages')}")
                    st.text_area(
                        "前 1000 字预览",
                        value=str(extraction_report.get("preview", "")),
                        height=220,
                        disabled=True,
                    )
                    total_pages = max(1, int(extraction_report.get("total_pages", 1)))
                    char_count = int(extraction_report.get("char_count", 0))
                    if char_count < 800 or char_count / total_pages < 120:
                        st.warning("PDF 可能是扫描版或图片型 PDF，当前无法完整识别，请使用可复制文字的 PDF 或复制粘贴文本。")
        elif file_suffix == ".pdf":
            st.warning("没有从这个 PDF 中提取到可选中文本。如果这是扫描版 PDF，需要先做 OCR。")
            if extraction_report:
                with st.expander("PDF 提取报告", expanded=True):
                    st.markdown(f"**文件名：** {extraction_report.get('filename', '')}")
                    st.markdown(f"**总页数：** {extraction_report.get('total_pages', 0)}")
                    st.markdown(f"**成功提取页数：** {extraction_report.get('successful_pages', 0)}")
                    st.markdown(f"**提取字符数：** {extraction_report.get('char_count', 0)}")
                    st.warning("PDF 可能是扫描版或图片型 PDF，当前无法完整识别，请使用可复制文字的 PDF 或复制粘贴文本。")
        elif file_suffix == ".docx":
            st.warning("没有从这个 Word 文档中提取到文本。请确认文档不是空白或受保护文件。")
        elif file_suffix == ".pptx":
            st.warning("没有从这个 PowerPoint 中提取到文本。请确认幻灯片不是纯图片，或复制粘贴文字内容。")
        elif file_suffix == ".txt":
            st.warning("没有从这个 TXT 文件中提取到文本。请确认文件不是空白。")
    except Exception as exc:
        st.error(f"无法提取这个课程文档的文本：{exc}")

text = st.text_area(
    t("course_content"),
    value=uploaded_text or sample,
    height=220,
    placeholder=t("course_placeholder"),
)
st.session_state["last_course_text"] = text

subject = st.selectbox(
    t("subject"),
    SUBJECT_VALUES,
    format_func=format_subject_option,
    key="study_pack_subject",
)
if not subject:
    subject = "Dentistry"

generation_depth = st.selectbox(
    t("generation_depth"),
    ["快速总结", "标准复习包", "考前冲刺包", "详细逐题版"],
    index=1,
    format_func={
        "快速总结": t("depth_quick"),
        "标准复习包": t("depth_standard"),
        "考前冲刺包": t("depth_cram"),
        "详细逐题版": t("depth_detailed"),
    }.get,
    help=t("generation_depth_help"),
)
st.caption(t("generation_depth_caption"))

expected_section_count = st.number_input(
    t("expected_section_count"),
    min_value=0,
    max_value=200,
    value=0,
    step=1,
    help=t("expected_section_help"),
)

col_a, col_b = st.columns([1.25, 3.75], vertical_alignment="center")
with col_a:
    generate = st.button(t("generate_study_pack"), type="primary", use_container_width=True)
with col_b:
    st.caption(t("small_test_tip"))

st.markdown("</div>", unsafe_allow_html=True)

st.markdown(f"### {t('recent_study_pack_records')}")
if st.session_state.get("study_pack_records_error"):
    st.error(f"读取复习包记录失败：{st.session_state['study_pack_records_error']}")
recent_study_pack_records = st.session_state.get("study_pack_records", [])
if not recent_study_pack_records:
    st.info(t("no_study_pack_records"))
else:
    for index, record in enumerate(recent_study_pack_records[:5], start=1):
        created_at = str(record.get("created_at", ""))[:19].replace("T", " ")
        label = f"{index}. {record.get('source_title') or record.get('subject') or 'Study Pack'} · {created_at}"
        with st.expander(label):
            st.caption(f"科目：{record.get('subject', 'Dentistry')}")
            source_preview = str(record.get("source_text") or "")[:500]
            if source_preview:
                st.write(source_preview)
            if st.button(t("open_study_pack"), key=f"open_study_pack_{record.get('id')}"):
                reopened_pack = record.get("generated_pack") or {}
                reopened_subject = record.get("subject") or "Dentistry"
                st.session_state["study_pack_result"] = reopened_pack
                st.session_state["study_pack_md_bytes"] = str(record.get("markdown_export") or "").encode("utf-8")
                try:
                    st.session_state["study_pack_pdf_bytes"] = build_study_pack_pdf(reopened_pack, reopened_subject)
                except Exception:
                    st.session_state["study_pack_pdf_bytes"] = None
                st.session_state["anki_csv_bytes"] = build_anki_csv_bytes(reopened_pack, reopened_subject)
                st.success("已打开保存的复习包。")
                st.rerun()


if generate:
    if not text.strip():
        st.error("请先粘贴英文课程内容，或上传 PDF / Word(docx) / TXT 文档。")
        st.stop()

    pack = run_with_dental_loader(
        "正在按题号/章节生成复习包，大文件可能需要几分钟..." if get_ui_lang() == "en" else "正在按题号/章节生成复习包，大文件可能需要几分钟...",
        lambda: generate_study_pack(
            text,
            subject,
            generation_depth,
            int(expected_section_count) or None,
            source_filename=uploaded_source_name or (uploaded_course_file.name if uploaded_course_file is not None else "Pasted course text"),
            page_or_slide_count=uploaded_page_or_slide_count,
        ),
        "课程内容会分块处理，请不要关闭页面。" if get_ui_lang() == "en" else "课程内容会分块处理，请不要关闭页面。",
    )

    st.markdown('<div class="result-wrap">', unsafe_allow_html=True)
    st.markdown("### 复习包")
    status_message = pack.get("status_message")
    if pack.get("mode") == "fallback":
        st.warning(status_message or "当前使用本地备用模式。")
    elif status_message:
        st.success(status_message)

    pdf_bytes = None
    pdf_error = None
    try:
        pdf_bytes = build_study_pack_pdf(pack, subject)
    except Exception as exc:
        pdf_error = str(exc)

    markdown_bytes = build_study_pack_markdown(pack, subject)
    anki_csv_bytes = build_anki_csv_bytes(pack, subject)
    st.session_state["study_pack_result"] = pack
    st.session_state["study_pack_pdf_bytes"] = pdf_bytes
    st.session_state["study_pack_md_bytes"] = markdown_bytes
    st.session_state["anki_csv_bytes"] = anki_csv_bytes
    try:
        source_title = uploaded_course_file.name if uploaded_course_file is not None else "Pasted course text"
        saved_record = save_study_pack_record(
            subject,
            source_title,
            text,
            pack,
            markdown_bytes.decode("utf-8", errors="ignore"),
        )
        if saved_record:
            st.session_state.setdefault("study_pack_records", [])
            st.session_state["study_pack_records"].insert(0, saved_record)
            st.session_state["study_pack_records"] = st.session_state["study_pack_records"][:10]
        increment_usage("study_pack", 1)
    except Exception as exc:
        st.warning(f"复习包记录保存失败：{exc}")

    if pdf_bytes:
        st.download_button(
            label=t("download_pdf_pack"),
            data=pdf_bytes,
            file_name="dentpilot_study_pack.pdf",
            mime="application/pdf",
            use_container_width=True,
        )
    else:
        st.warning(f"PDF 下载失败时，可先下载 Markdown 版本。{pdf_error or ''}")

    st.download_button(
        label=t("download_md_pack"),
        data=markdown_bytes,
        file_name="dentpilot_study_pack.md",
        mime="text/markdown",
        use_container_width=True,
    )
    st.info("手机端点击下载后，请在浏览器下载管理或文件 App 中查看。如果 PDF 无法打开，请先下载 Markdown 版本。")
    modules = pack.get("study_modules", [])
    coverage = pack.get("coverage_report", {})
    processing_details = pack.get("processing_details", {})
    detected_sections = coverage.get("detected_sections", len(modules))
    original_detected_sections = coverage.get("original_detected_sections", detected_sections)
    generated_sections = coverage.get("generated_sections", len(modules))
    coverage_percent = coverage.get("coverage_percent", 100 if modules else 0)

    with st.expander("处理详情 / Processing Details", expanded=True):
        detail_col_1, detail_col_2, detail_col_3 = st.columns(3)
        detail_col_1.metric("提取字符数", processing_details.get("extracted_char_count", len(text)))
        detail_col_2.metric("分块数量", processing_details.get("chunk_count", detected_sections))
        detail_col_3.metric("已处理分块", processing_details.get("processed_chunk_count", generated_sections))
        st.markdown(f"**文件名：** {processing_details.get('uploaded_file_name') or (uploaded_course_file.name if uploaded_course_file is not None else 'Pasted course text')}")
        st.markdown(f"**估计页数/幻灯片数：** {processing_details.get('estimated_page_or_slide_count', uploaded_page_or_slide_count)}")
        st.markdown(f"**Chunk size：** {processing_details.get('chunk_size', '')}")
        st.markdown(f"**Overlap size：** {processing_details.get('overlap_size', '')}")
        st.markdown(f"**失败分块数：** {processing_details.get('failed_chunk_count', 0)}")
        st.markdown(f"**备用分块数：** {processing_details.get('fallback_chunk_count', 0)}")
        st.markdown(f"**跳过分块数：** {processing_details.get('skipped_chunk_count', 0)}")
        st.markdown(f"**模型：** {processing_details.get('model_used', '')}")

        preview_col_1, preview_col_2, preview_col_3 = st.columns(3)
        preview_col_1.text_area("第一个 chunk 前 200 字", processing_details.get("first_chunk_preview", ""), height=160, disabled=True)
        preview_col_2.text_area("中间 chunk 前 200 字", processing_details.get("middle_chunk_preview", ""), height=160, disabled=True)
        preview_col_3.text_area("最后 chunk 前 200 字", processing_details.get("last_chunk_preview", ""), height=160, disabled=True)

        st.markdown("**早期 chunks 主题：** " + "；".join(processing_details.get("early_topics", [])))
        st.markdown("**中间 chunks 主题：** " + "；".join(processing_details.get("middle_topics", [])))
        st.markdown("**后期 chunks 主题：** " + "；".join(processing_details.get("late_topics", [])))
        chunk_statuses = processing_details.get("chunk_statuses", [])
        if chunk_statuses:
            st.dataframe(pd.DataFrame(chunk_statuses), use_container_width=True, hide_index=True)
        if processing_details.get("failed_chunk_count", 0) or processing_details.get("fallback_chunk_count", 0):
            st.warning("以下分块处理失败或未充分覆盖，请查看上方状态表；系统已尽量用备用模块保留这些分块。")

    if modules:
        st.info(
            f"检测到 {original_detected_sections} 个原始考试题/章节，拆分为 {detected_sections} 个处理分块，"
            f"已生成 {generated_sections} 个复习模块，"
            f"覆盖率：{coverage_percent}%。"
        )
        if coverage.get("has_missing"):
            st.warning("部分章节未生成，请使用详细逐题版或减少上传内容。")
        if coverage.get("expected_mismatch"):
            st.error(
                f"你填写的预期题目数量是 {coverage.get('expected_sections')}，"
                f"但系统只检测到 {detected_sections} 个。请检查原文提取结果，或把文件另存为 DOCX/TXT 后重试。"
            )

    tab_overview, tab_modules, tab_terms, tab_oral, tab_quiz, tab_anki, tab_summary, tab_coverage = st.tabs(
        ["总览", "逐题讲解", "高频术语", "口试题库", "Quiz", "Anki", "考前总结", "覆盖率检查"]
    )

    with tab_overview:
        st.subheader("总览")
        st.write(pack.get("exam_summary", ""))
        if modules:
            overview_rows = [
                {
                    "题号/章节": module.get("section_number"),
                    "标题": module.get("title"),
                    "Must-know 数": len(module.get("must_know", [])),
                    "Quiz 数": len(module.get("quiz", [])),
                    "Anki 数": len(module.get("flashcards", [])),
                }
                for module in modules
            ]
            st.dataframe(pd.DataFrame(overview_rows), use_container_width=True)
        else:
            st.write(pack.get("chinese_explanation", ""))
            st.subheader("核心概念")
            for item in pack.get("key_concepts", []):
                st.markdown(f"- {item}")

    with tab_modules:
        st.subheader("逐题讲解")
        if modules:
            for module in modules:
                title_text = f"Question {module.get('section_number')}: {module.get('title', '')}"
                with st.expander(title_text, expanded=False):
                    st.markdown("**中文核心讲解**")
                    st.write(module.get("chinese_core_explanation", ""))
                    st.markdown("**Must-know points**")
                    for item in module.get("must_know", []):
                        st.markdown(f"- {item}")
                    st.markdown("**Common mistakes**")
                    for item in module.get("common_mistakes", []):
                        st.markdown(f"- {item}")
                    st.markdown("**Short answer template**")
                    st.write(module.get("short_answer_template", ""))
                    if module.get("source_excerpt"):
                        with st.expander("原文片段"):
                            st.write(module.get("source_excerpt"))
        else:
            st.write(pack.get("chinese_explanation", ""))

    with tab_terms:
        st.subheader("高频术语")
        if pack.get("glossary"):
            df_terms = pd.DataFrame(pack["glossary"])
            st.dataframe(df_terms, use_container_width=True)
            csv_terms = df_terms.to_csv(index=False).encode("utf-8-sig")
            st.download_button("下载术语表 CSV", csv_terms, "glossary.csv", "text/csv")
        else:
            st.info("没有生成术语表。")

    with tab_oral:
        st.subheader("口试题库")
        if modules:
            for module in modules:
                with st.expander(f"Question {module.get('section_number')}: {module.get('title', '')}"):
                    st.markdown("**Likely oral exam questions**")
                    for item in module.get("oral_exam_questions", []):
                        st.markdown(f"- {item}")
                    st.markdown("**Follow-up questions**")
                    for item in module.get("follow_up_questions", []):
                        st.markdown(f"- {item}")
        else:
            for q in pack.get("quiz", []):
                st.markdown(f"- {q.get('question', '')}")

    with tab_quiz:
        st.subheader("Quiz")
        if modules:
            for module in modules:
                with st.expander(f"Question {module.get('section_number')}: {module.get('title', '')}"):
                    for i, q in enumerate(module.get("quiz", []), start=1):
                        st.markdown(f"**{i}. [{q.get('question_type', 'quiz')}] {q.get('question', '')}**")
                        for option in q.get("options", []):
                            st.write(f"- {option}")
                        with st.expander("查看答案与中文解析"):
                            st.write(f"答案：{q.get('answer', '')}")
                            st.write(q.get("explanation_zh", ""))
        else:
            for i, q in enumerate(pack.get("quiz", []), start=1):
                st.markdown(f"**Q{i}. {q.get('question', '')}**")
                for option in q.get("options", []):
                    st.write(f"- {option}")
                    with st.expander("查看答案与中文解析"):
                        st.write(f"答案：{q.get('answer', '')}")
                    st.write(q.get("explanation_zh", ""))

    with tab_anki:
        st.subheader("Anki 卡片")
        df_cards = pd.DataFrame(pack.get("flashcards", []))
        st.dataframe(df_cards, use_container_width=True)

        st.download_button(
            label=t("download_anki_csv"),
            data=st.session_state.get("anki_csv_bytes") or build_anki_csv_bytes(pack, subject),
            file_name="medstudy_anki_cards.csv",
            mime="text/csv",
            use_container_width=True,
        )
    with tab_summary:
        st.subheader("考前总结")
        st.write(pack.get("exam_summary", ""))
        if modules:
            st.markdown("**考前冲刺顺序**")
            st.write("1. 先背每题 Must-know points；2. 再看 Common mistakes；3. 最后用口试题库自测英文表达。")
        st.info("PDF 导出已包含逐题讲解、口试题、易错点、Anki 卡片和覆盖率报告。")

    with tab_coverage:
        st.subheader("覆盖率检查")
        st.metric("检测到的原始问题/章节", original_detected_sections)
        st.metric("处理分块数", detected_sections)
        st.metric("实际生成的复习模块", generated_sections)
        st.metric("覆盖率", f"{coverage_percent}%")
        if coverage.get("expected_sections"):
            st.metric("预期题目数量", coverage.get("expected_sections"))
        if coverage.get("expected_mismatch"):
            st.error(
                "检测数量和预期题目数量不一致。建议先查看下方“检测到的题目/章节”，"
                "确认缺的是哪些题；如果 PDF 提取不完整，请改用 DOCX 或 TXT。"
            )
        if coverage.get("missing_sections"):
            st.warning(f"遗漏章节：{coverage.get('missing_sections')}")
        elif modules:
            st.success("没有检测到遗漏。")
        detected_titles = coverage.get("detected_titles", [])
        if detected_titles:
            st.write("检测到的题目/章节：")
            st.dataframe(pd.DataFrame(detected_titles), use_container_width=True)
        st.write(f"检测方式：{coverage.get('detection_method', 'legacy')}")
    st.markdown("</div>", unsafe_allow_html=True)
else:
    cached_pdf_bytes = st.session_state.get("study_pack_pdf_bytes")
    cached_md_bytes = st.session_state.get("study_pack_md_bytes")
    cached_anki_bytes = st.session_state.get("anki_csv_bytes")
    if cached_pdf_bytes or cached_md_bytes or cached_anki_bytes:
        st.markdown('<div class="result-wrap">', unsafe_allow_html=True)
        st.markdown("### 已生成的下载文件")
        if cached_pdf_bytes:
            st.download_button(
                label=t("download_pdf_pack"),
                data=cached_pdf_bytes,
                file_name="dentpilot_study_pack.pdf",
                mime="application/pdf",
                use_container_width=True,
            )
        if cached_md_bytes:
            st.download_button(
                label=t("download_md_pack"),
                data=cached_md_bytes,
                file_name="dentpilot_study_pack.md",
                mime="text/markdown",
                use_container_width=True,
            )
        if cached_anki_bytes:
            st.download_button(
                label=t("download_anki_csv"),
                data=cached_anki_bytes,
                file_name="medstudy_anki_cards.csv",
                mime="text/csv",
                use_container_width=True,
            )
        st.info("手机端点击下载后，请在浏览器下载管理或文件 App 中查看。如果 PDF 无法打开，请先下载 Markdown 版本。")
        st.markdown("</div>", unsafe_allow_html=True)
    else:
        st.info("点击“生成复习包”开始。")

st.markdown(
    """
    <div class="footer">
        DentPilot AI 帮助英授牙科/医学学生把英文课程资料整理成中文讲解、术语复习、Quiz 自测、Anki 卡片和 PDF 复习包。
    </div>
    """,
    unsafe_allow_html=True,
)
